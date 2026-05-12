"""Build the Enterprise Agentic Stack layer-cake deck — a white-background
exec-grade architecture slide for technical-but-business audiences.

Output: docs/decks/enterprise-stack.pptx

Two slides:
  1. Cover — title + thesis (white background)
  2. The expanded layer-cake architecture — 5 stacked layers + an AIOps
     foundation band that spans Microsoft + Oracle infra agents.

White background is deliberate — this slide is designed to drop into
existing exec presentations that use light slide masters.

Layer cake (top to bottom):
  1. Agentic Business Process Layer — Microsoft Copilot + AI applications
  2. Intelligence Layer — Fabric IQ · Foundry IQ · Work IQ
  3. Implementation Layer — Microsoft AI + cloud services that power this
  4. Agentic Integration to Oracle Data — Fabric Mirror / MCP / A2A lanes
  5. Data Layer — Oracle data of record (Fusion, OD@A, OCI, on-prem)
  Foundation: AIOps — Microsoft infra agents | Oracle infra agents
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette — recoloured for white background readability
# ---------------------------------------------------------------------------
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY_900 = RGBColor(0x0B, 0x12, 0x20)
NAVY_700 = RGBColor(0x1A, 0x25, 0x40)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_300 = RGBColor(0xC2, 0xC9, 0xD4)
SLATE_100 = RGBColor(0xEF, 0xF1, 0xF5)

ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)
ORACLE_RED_DK = RGBColor(0x8A, 0x2D, 0x1E)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
AZURE_BLUE_DK = RGBColor(0x00, 0x4F, 0x8B)
IQ_YELLOW = RGBColor(0xE0, 0xA8, 0x00)  # Slightly darker yellow for white-bg readability
IQ_YELLOW_DK = RGBColor(0x9A, 0x73, 0x00)
IQ_TEAL = RGBColor(0x0F, 0xA8, 0x9B)
IQ_TEAL_DK = RGBColor(0x0C, 0x70, 0x68)
PURPLE = RGBColor(0x6A, 0x4C, 0xC0)  # Bridge layer color
PURPLE_DK = RGBColor(0x44, 0x2E, 0x82)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_white_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.shadow.inherit = False


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=NAVY_900, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_layer_card(slide, left, top, width, height, border_color, *,
                    tint_alpha=0.88):
    """Layer card with a brand-coloured border and a very light tint fill."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    card.fill.solid()
    card.fill.fore_color.rgb = border_color
    card.fill.transparency = tint_alpha
    card.shadow.inherit = False
    return card


def add_pill(slide, left, top, width, height, text, fill_color,
             text_color, *, size=9, bold=True, border_color=None):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    if border_color is None:
        pill.line.color.rgb = fill_color
    else:
        pill.line.color.rgb = border_color
    pill.line.width = Pt(0.75)
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill_color
    pill.shadow.inherit = False
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = "Consolas"


def add_filled_block(slide, left, top, width, height, fill, line_color, *,
                     transparency=0.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.10
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if transparency > 0:
        shape.fill.transparency = transparency
    shape.shadow.inherit = False
    return shape


def add_circle_badge(slide, cx, cy, radius, text, fill, text_color):
    diameter = 2 * radius
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     cx - radius, cy - radius, diameter, diameter)
    circle.line.color.rgb = fill
    circle.line.width = Pt(1)
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Consolas"


# ---------------------------------------------------------------------------
# Slide 1 — Cover (white background)
# ---------------------------------------------------------------------------
def cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)

    # Top accent strip (thin gradient bar across the top)
    accent_segments = [
        (Inches(0), IQ_TEAL),
        (Inches(2.67), IQ_YELLOW),
        (Inches(5.33), AZURE_BLUE),
        (Inches(8.0), PURPLE),
        (Inches(10.67), ORACLE_RED),
    ]
    for x, color in accent_segments:
        seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(0), Inches(2.67), Inches(0.12))
        seg.line.fill.background()
        seg.fill.solid()
        seg.fill.fore_color.rgb = color
        seg.shadow.inherit = False

    # Eyebrow
    add_text(
        slide, Inches(0.6), Inches(0.55), Inches(12), Inches(0.4),
        "ENTERPRISE AGENTIC STACK  ·  ORACLE × MICROSOFT  ·  END-TO-END",
        size=11, color=SLATE_500, font="Consolas",
    )

    # Title
    add_text(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(3.2),
        "The agentic stack from\nworker to data — layer by layer.",
        size=42, bold=True, color=NAVY_900,
    )

    # Subtitle
    add_text(
        slide, Inches(0.6), Inches(4.4), Inches(12), Inches(2.0),
        "Five layers and one AIOps foundation that together make Azure the best place to run Oracle workloads.\n"
        "Microsoft Copilot at the worker surface. The IQ layer in the middle. Oracle data of record at the base. "
        "AIOps agents — distributed across Microsoft infra and Oracle infra — power the whole stack.",
        size=17, color=SLATE_700,
    )

    # Bottom strip
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4),
        prs.slide_width, Inches(0.4),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = SLATE_100
    band.shadow.inherit = False
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.35), Inches(12), Inches(0.32),
        "Designed to drop into exec presentations on light slide masters.  ·  See companion slide for the full layer cake.",
        size=10, color=SLATE_500, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 2 — The expanded layer-cake architecture (the show-piece)
# ---------------------------------------------------------------------------
def layer_cake_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(slide, prs)

    # ----- Title strip -----
    add_text(
        slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.32),
        "ENTERPRISE AGENTIC STACK  ·  LAYER CAKE FROM WORKER TO DATA",
        size=11, color=SLATE_500, font="Consolas",
    )
    add_text(
        slide, Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.55),
        "The complete end-to-end stack — Oracle data, Microsoft agents.",
        size=22, bold=True, color=NAVY_900,
    )

    # ----- Layout constants -----
    CONTENT_LEFT = Inches(0.5)
    CONTENT_W = Inches(12.3)
    LAYER_TOP = Inches(1.3)
    LAYER_H = Inches(0.95)
    LAYER_GAP = Inches(0.08)

    # Layers definition — order top to bottom
    # Each: (eyebrow, title, pills [(text, fill, text_color, border)], border_color, badge_label, badge_text_color)
    layers = [
        {
            "eyebrow": "LAYER 1 · AGENTIC BUSINESS PROCESS",
            "title": "Microsoft Copilot + AI applications",
            "border": IQ_TEAL,
            "border_dk": IQ_TEAL_DK,
            "badge": "1",
            "pills": [
                "Microsoft Copilot",
                "Custom AI apps",
                "Industry workflows",
                "Field-worker apps",
                "Copilot extensions",
            ],
        },
        {
            "eyebrow": "LAYER 2 · INTELLIGENCE",
            "title": "The IQ layer — Fabric IQ · Foundry IQ · Work IQ",
            "border": IQ_YELLOW,
            "border_dk": IQ_YELLOW_DK,
            "badge": "2",
            "pills": [
                "Fabric IQ · semantic",
                "Foundry IQ · reasoning",
                "Work IQ · M365 Graph",
                "Ontology Bridge",
                "Governance Bridge",
            ],
        },
        {
            "eyebrow": "LAYER 3 · IMPLEMENTATION",
            "title": "Microsoft AI + cloud services that power this",
            "border": AZURE_BLUE,
            "border_dk": AZURE_BLUE_DK,
            "badge": "3",
            "pills": [
                "Microsoft Agent Framework",
                "Azure AI Foundry",
                "Microsoft Agent 365",
                "Copilot Studio",
                "Microsoft Fabric",
                "Microsoft Purview",
                "Azure OpenAI",
            ],
        },
        {
            "eyebrow": "LAYER 4 · AGENTIC INTEGRATION TO ORACLE DATA",
            "title": "Three lanes between Microsoft and Oracle agents",
            "border": PURPLE,
            "border_dk": PURPLE_DK,
            "badge": "4",
            "pills": None,  # rendered as 3 lanes — see special handling below
        },
        {
            "eyebrow": "LAYER 5 · DATA",
            "title": "Oracle data of record — wherever it lives",
            "border": ORACLE_RED,
            "border_dk": ORACLE_RED_DK,
            "badge": "5",
            "pills": [
                "Oracle Fusion (SaaS)",
                "Oracle Database@Azure",
                "Oracle Cloud Infrastructure",
                "On-prem databases",
            ],
        },
    ]

    # ----- Render 5 layers stacked -----
    layer_label_w = Inches(3.4)  # left strip for badge + label
    pill_area_left_offset = layer_label_w + Inches(0.15)

    for i, layer in enumerate(layers):
        y = LAYER_TOP + i * (LAYER_H + LAYER_GAP)

        # Layer card with light tint
        add_layer_card(slide, CONTENT_LEFT, y, CONTENT_W, LAYER_H,
                       border_color=layer["border"], tint_alpha=0.92)

        # Number badge (circle) on the left
        badge_cx = CONTENT_LEFT + Inches(0.45)
        badge_cy = y + LAYER_H / 2
        add_circle_badge(slide, badge_cx, badge_cy, Inches(0.22),
                         layer["badge"], fill=layer["border"], text_color=WHITE)

        # Eyebrow + title
        add_text(slide, CONTENT_LEFT + Inches(0.85), y + Inches(0.12),
                 layer_label_w - Inches(0.5), Inches(0.28),
                 layer["eyebrow"], size=9.5, bold=True,
                 color=layer["border_dk"], font="Consolas")
        add_text(slide, CONTENT_LEFT + Inches(0.85), y + Inches(0.4),
                 layer_label_w - Inches(0.5), Inches(0.55),
                 layer["title"], size=13.5, bold=True, color=NAVY_900)

        # Pill area on the right
        if layer["pills"] is not None:
            pills = layer["pills"]
            pill_area_x = CONTENT_LEFT + pill_area_left_offset
            pill_area_w = CONTENT_W - pill_area_left_offset - Inches(0.2)
            # Compute pill widths from text length, but normalised to fit
            pill_h = Inches(0.36)
            pill_y = y + (LAYER_H - pill_h) / 2
            gap = Inches(0.1)

            # Estimate widths from char count
            raw_widths = [Inches(0.18 + 0.085 * len(p)) for p in pills]
            total = sum(w for w in raw_widths) + gap * (len(pills) - 1)
            scale = min(1.0, float(pill_area_w) / float(total))
            scaled_widths = [Inches(float(w) * scale / 914400) for w in raw_widths]

            pill_x = pill_area_x
            for j, text in enumerate(pills):
                w = scaled_widths[j]
                # Pill — solid brand-color fill, white text
                add_pill(slide, pill_x, pill_y, w, pill_h, text,
                         fill_color=layer["border"], text_color=WHITE, size=9)
                pill_x += w + gap
        else:
            # Layer 4 — three lanes (Fabric Mirror / MCP / A2A) rendered as
            # three sub-cards inside the pill area.
            pill_area_x = CONTENT_LEFT + pill_area_left_offset
            pill_area_w = CONTENT_W - pill_area_left_offset - Inches(0.2)
            lane_h = Inches(0.65)
            lane_y = y + (LAYER_H - lane_h) / 2
            lane_gap = Inches(0.1)
            lane_w = (pill_area_w - 2 * lane_gap) / 3

            lanes = [
                ("Fabric Mirror → OneLake", "ORACLE → MS  ·  data plane", AZURE_BLUE),
                ("MCP tools", "MS → ORACLE  ·  tool handoff", IQ_TEAL),
                ("A2A handoff", "BIDIRECTIONAL  ·  peer agents", ORACLE_RED),
            ]
            lx = pill_area_x
            for label, sub, color in lanes:
                lane = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               lx, lane_y, lane_w, lane_h)
                lane.adjustments[0] = 0.15
                lane.line.color.rgb = color
                lane.line.width = Pt(1.5)
                lane.fill.solid()
                lane.fill.fore_color.rgb = color
                lane.fill.transparency = 0.85
                lane.shadow.inherit = False
                add_text(slide, lx, lane_y + Inches(0.07),
                         lane_w, Inches(0.28),
                         label, size=11, bold=True, color=NAVY_900,
                         align=PP_ALIGN.CENTER)
                add_text(slide, lx, lane_y + Inches(0.36),
                         lane_w, Inches(0.25),
                         sub, size=8.5, color=color,
                         font="Consolas", align=PP_ALIGN.CENTER, bold=True)
                lx += lane_w + lane_gap

    # ----- AIOps foundation band (spans Microsoft + Oracle infra agents) -----
    foundation_top = LAYER_TOP + 5 * (LAYER_H + LAYER_GAP) + Inches(0.05)
    foundation_h = Inches(0.95)

    # Outer card
    add_layer_card(slide, CONTENT_LEFT, foundation_top, CONTENT_W, foundation_h,
                   border_color=SLATE_500, tint_alpha=0.94)

    # Foundation badge
    badge_cx = CONTENT_LEFT + Inches(0.45)
    badge_cy = foundation_top + foundation_h / 2
    add_circle_badge(slide, badge_cx, badge_cy, Inches(0.22),
                     "∞", fill=NAVY_900, text_color=WHITE)

    # Eyebrow + title
    add_text(slide, CONTENT_LEFT + Inches(0.85), foundation_top + Inches(0.12),
             Inches(3), Inches(0.28),
             "FOUNDATION · AIOps", size=9.5, bold=True, color=NAVY_700, font="Consolas")
    add_text(slide, CONTENT_LEFT + Inches(0.85), foundation_top + Inches(0.4),
             Inches(3), Inches(0.55),
             "Powers the entire stack",
             size=13.5, bold=True, color=NAVY_900)

    # Split: Microsoft infra agents | Oracle infra agents
    split_x = CONTENT_LEFT + pill_area_left_offset
    split_w = CONTENT_W - pill_area_left_offset - Inches(0.2)
    half_w = (split_w - Inches(0.15)) / 2
    split_y = foundation_top + Inches(0.12)
    split_h = foundation_h - Inches(0.24)

    # Left: Microsoft infra agents
    ms = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 split_x, split_y, half_w, split_h)
    ms.adjustments[0] = 0.12
    ms.line.color.rgb = AZURE_BLUE
    ms.line.width = Pt(1.5)
    ms.fill.solid()
    ms.fill.fore_color.rgb = AZURE_BLUE
    ms.fill.transparency = 0.82
    ms.shadow.inherit = False
    add_text(slide, split_x + Inches(0.15), split_y + Inches(0.06),
             half_w - Inches(0.3), Inches(0.25),
             "MICROSOFT INFRA AGENTS", size=9, bold=True,
             color=AZURE_BLUE_DK, font="Consolas")
    add_text(slide, split_x + Inches(0.15), split_y + Inches(0.32),
             half_w - Inches(0.3), Inches(0.35),
             "Config Drift  ·  Troubleshooting & RCA  ·  Workload Health  ·  Configuration Validation  ·  WVI-level rollups",
             size=10.5, color=NAVY_900)

    # Right: Oracle infra agents
    ora = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  split_x + half_w + Inches(0.15), split_y,
                                  half_w, split_h)
    ora.adjustments[0] = 0.12
    ora.line.color.rgb = ORACLE_RED
    ora.line.width = Pt(1.5)
    ora.fill.solid()
    ora.fill.fore_color.rgb = ORACLE_RED
    ora.fill.transparency = 0.82
    ora.shadow.inherit = False
    add_text(slide, split_x + half_w + Inches(0.3), split_y + Inches(0.06),
             half_w - Inches(0.3), Inches(0.25),
             "ORACLE INFRA AGENTS", size=9, bold=True,
             color=ORACLE_RED_DK, font="Consolas")
    add_text(slide, split_x + half_w + Inches(0.3), split_y + Inches(0.32),
             half_w - Inches(0.3), Inches(0.35),
             "Oracle Audit Vault  ·  Data Safe drift  ·  Database AI agents  ·  Fusion AI Agents  ·  Cloud Lift",
             size=10.5, color=NAVY_900)

    # ----- Bottom tagline -----
    tagline_y = foundation_top + foundation_h + Inches(0.08)
    add_text(slide, CONTENT_LEFT, tagline_y, CONTENT_W, Inches(0.3),
             "Make Azure the best place to run Oracle workloads.",
             size=12, bold=True, color=NAVY_900, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    cover_slide(prs)
    layer_cake_slide(prs)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "decks"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "enterprise-stack.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build()
