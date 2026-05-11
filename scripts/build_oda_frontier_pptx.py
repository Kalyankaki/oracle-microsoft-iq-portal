"""Build the OD@A Frontier exec deck — the Microsoft-internal briefing deck
that mirrors the /oda-frontier portal page.

Output: docs/decks/oda-frontier.pptx

Seven slides:
  1. Title — "OD@A is where every Oracle workload's modernisation lands."
  2. Why OD@A is the anchor — 4 anchor properties (2x2 grid)
  3. The customer journey — 5 stages with OD@A as the gravitational anchor
  4. The cloud-path matrix — 4 destinations, OD@A emphasised
  5. Technical architecture — 5-layer summary (estate -> OD@A -> IQ stack)
  6. The consumption story — per-stage Microsoft + Oracle revenue lines
  7. The exec ask — 4 asks for Microsoft leadership + the bet

Visual conventions, brand palette, and helper functions mirror
scripts/build_scenarios_pptx.py so the two decks read as a family.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette (matches the portal Tailwind config + scenarios deck)
# ---------------------------------------------------------------------------
NAVY_900 = RGBColor(0x0B, 0x12, 0x20)
NAVY_700 = RGBColor(0x1A, 0x25, 0x40)
ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
IQ_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
IQ_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_85 = RGBColor(0xD9, 0xDD, 0xE5)
WHITE_70 = RGBColor(0xB0, 0xB7, 0xC4)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


# ---------------------------------------------------------------------------
# Helpers (same shapes as the scenarios deck, kept inline so the script
# stands on its own — see scripts/build_scenarios_pptx.py for the precedent)
# ---------------------------------------------------------------------------
def add_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_900
    bg.shadow.inherit = False


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_pill(slide, left, top, width, height, text, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(0.75)
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.fill.transparency = 0.85
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Consolas"


def add_card(slide, left, top, width, height, *, border=MUTED, fill=NAVY_700, fill_alpha=0.6):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.line.color.rgb = border
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.fill.transparency = fill_alpha
    card.shadow.inherit = False
    return card


def add_footer(slide, prs, text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4),
        prs.slide_width, Inches(0.4),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.36), Inches(12), Inches(0.32),
        text, size=9, color=MUTED, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
        "OD@A  ·  ANCHOR FOR FRONTIER TRANSFORMATION  ·  ORACLE × MICROSOFT",
        size=11, color=ORACLE_RED, font="Consolas",
    )

    add_text(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(2.6),
        "Oracle Database@Azure\nis where every Oracle workload's\nmodernisation lands.",
        size=42, bold=True, color=WHITE,
    )

    add_text(
        slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.0),
        "The agentic alliance compounds value from a stable anchor — and OD@A is the anchor.\n"
        "This deck walks the end-to-end customer journey, the cloud-path decision tree, the technical "
        "stack, the consumption story, and what we need from Microsoft leadership to make OD@A the "
        "default destination for Oracle workloads.",
        size=17, color=WHITE_85,
    )

    # Footer band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5),
        prs.slide_width, Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.45), Inches(12), Inches(0.4),
        "Oracle anchors data of record on OD@A.  ·  Microsoft extends the worker surface through Copilot.  ·  Together: frontier transformation.",
        size=10, color=MUTED, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 2 — Why OD@A is the anchor (4 cards in 2x2 grid)
# ---------------------------------------------------------------------------
ANCHOR_CLAIMS = [
    (
        "Lowest-latency Fabric mirror",
        "Native Fabric Mirroring for OD@A (Ignite '25) gives the agent sub-second-fresh Oracle data without leaving Azure.",
        "No replication tax. No third-party connector. Oracle data behaves as if it were native Azure data.",
        AZURE_BLUE,
    ),
    (
        "Oracle-native features preserved",
        "Exadata, Autonomous Database, RAC, Data Guard — all run unchanged on OD@A.",
        "Same SQL, DBA workflows, backup/recovery, and licensing. Migration is move-and-modernise, not rewrite.",
        ORACLE_RED,
    ),
    (
        "No lift-and-shift mental tax",
        "Azure datacentres, jointly operated. One support call, one bill, one SLA.",
        "Arc + Entra + native peering make the seam between OD@A and the rest of Azure invisible to the worker.",
        IQ_YELLOW,
    ),
    (
        "Consumption flywheel scales with depth",
        "Every modernised workload lights up Exadata ratecards, GoldenGate, Fabric, Foundry, and Copilot — five revenue lines.",
        "OD@A compounds with Fabric, Foundry, and Copilot rather than competing with them.",
        IQ_TEAL,
    ),
]


def anchor_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "WHY OD@A IS THE ANCHOR",
        size=11, color=MUTED, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Four anchor properties — together, they make OD@A inevitable.",
        size=26, bold=True, color=WHITE,
    )

    # 2x2 grid of cards
    grid_top = Inches(1.85)
    card_w = Inches(6.0)
    card_h = Inches(2.55)
    gap = Inches(0.2)
    left_start = Inches(0.6)

    for i, (title, one_liner, proof, color) in enumerate(ANCHOR_CLAIMS):
        col = i % 2
        row = i // 2
        cl = left_start + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)

        add_card(slide, cl, ct, card_w, card_h, border=color, fill=NAVY_700)

        # Number chip in top-left
        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, cl + Inches(0.2), ct + Inches(0.2), Inches(0.35), Inches(0.35))
        chip.line.color.rgb = color
        chip.line.width = Pt(1)
        chip.fill.solid()
        chip.fill.fore_color.rgb = NAVY_900
        chip_tf = chip.text_frame
        chip_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        chip_tf.margin_left = chip_tf.margin_right = 0
        chip_tf.margin_top = chip_tf.margin_bottom = 0
        chip_p = chip_tf.paragraphs[0]
        chip_p.alignment = PP_ALIGN.CENTER
        chip_run = chip_p.add_run()
        chip_run.text = str(i + 1)
        chip_run.font.size = Pt(11)
        chip_run.font.bold = True
        chip_run.font.color.rgb = color
        chip_run.font.name = "Consolas"

        # Title
        add_text(slide, cl + Inches(0.7), ct + Inches(0.18), card_w - Inches(0.85), Inches(0.45),
                 title, size=15, bold=True, color=WHITE)
        # One-liner
        add_text(slide, cl + Inches(0.25), ct + Inches(0.75), card_w - Inches(0.4), Inches(0.95),
                 one_liner, size=12, color=WHITE_85)
        # Divider
        line = slide.shapes.add_connector(1, cl + Inches(0.25), ct + Inches(1.75),
                                          cl + card_w - Inches(0.25), ct + Inches(1.75))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        # Proof label
        add_text(slide, cl + Inches(0.25), ct + Inches(1.82), card_w - Inches(0.4), Inches(0.25),
                 "PROOF", size=9, color=MUTED, font="Consolas")
        add_text(slide, cl + Inches(0.25), ct + Inches(2.05), card_w - Inches(0.4), Inches(0.5),
                 proof, size=11, color=color, bold=True)

    add_footer(prs=prs, slide=slide,
               text="Reference: /oda-frontier  ·  Anchor properties · §2")


# ---------------------------------------------------------------------------
# Slide 3 — The customer journey (5 stages + OD@A anchor band)
# ---------------------------------------------------------------------------
JOURNEY_STAGES = [
    ("01", "Discover", "Inventory the on-prem Oracle estate — workloads, versions, dependencies, sensitive data, licence posture.", MUTED),
    ("02", "Plan", "Per-workload destination across four paths: OD@A, OCI, Azure DB for PostgreSQL, hybrid.", AZURE_BLUE),
    ("03", "Migrate", "ZDM-led migration · DMS re-platform · upgrade to compliant version on landing.", ORACLE_RED),
    ("04", "Modernise", "Apps move to AKS / App Service / Functions. Data of record stays anchored on OD@A.", IQ_YELLOW),
    ("05", "Activate AI", "Light up Fabric IQ + Foundry IQ + Work IQ on top. Agents grounded on Oracle data, surfaced in Copilot.", IQ_TEAL),
]


def journey_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE CUSTOMER JOURNEY",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Five stages. One anchor. Every workload converges on OD@A.",
        size=26, bold=True, color=WHITE,
    )

    # 5 stage cards horizontally
    stage_top = Inches(2.0)
    stage_h = Inches(3.0)
    stage_w = Inches(2.42)
    gap = Inches(0.1)
    left_start = Inches(0.6)

    for i, (num, label, blurb, color) in enumerate(JOURNEY_STAGES):
        cl = left_start + i * (stage_w + gap)
        add_card(slide, cl, stage_top, stage_w, stage_h, border=color, fill=NAVY_700)
        # Stage number
        add_text(slide, cl + Inches(0.2), stage_top + Inches(0.2), stage_w - Inches(0.4), Inches(0.4),
                 num, size=22, bold=True, color=color, font="Consolas")
        # Stage label
        add_text(slide, cl + Inches(0.2), stage_top + Inches(0.7), stage_w - Inches(0.4), Inches(0.45),
                 label, size=15, bold=True, color=WHITE)
        # Blurb
        add_text(slide, cl + Inches(0.2), stage_top + Inches(1.25), stage_w - Inches(0.4), Inches(1.7),
                 blurb, size=10.5, color=WHITE_85)
        # Bottom accent strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       cl, stage_top + stage_h - Inches(0.12),
                                       stage_w, Inches(0.12))
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = color

    # Anchor band at the bottom — OD@A
    anchor_top = Inches(5.4)
    anchor_left = Inches(2.4)
    anchor_w = Inches(8.5)
    anchor_h = Inches(1.3)
    add_card(slide, anchor_left, anchor_top, anchor_w, anchor_h,
             border=ORACLE_RED, fill=NAVY_700, fill_alpha=0.4)

    add_text(slide, anchor_left, anchor_top + Inches(0.18), anchor_w, Inches(0.45),
             "Oracle Database@Azure",
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, anchor_left, anchor_top + Inches(0.7), anchor_w, Inches(0.3),
             "Exadata · Autonomous · Database Service — Oracle on Azure datacentres",
             size=11.5, color=WHITE_85, align=PP_ALIGN.CENTER)
    add_text(slide, anchor_left, anchor_top + Inches(0.98), anchor_w, Inches(0.3),
             "ANCHOR  ·  EVERY STAGE CONVERGES HERE",
             size=10, color=ORACLE_RED, font="Consolas", align=PP_ALIGN.CENTER, bold=True)

    add_footer(slide=slide, prs=prs,
               text="Reference: /oda-frontier  ·  Customer journey · §3-4")


# ---------------------------------------------------------------------------
# Slide 4 — Cloud-path matrix (OD@A emphasised + 3 peer paths)
# ---------------------------------------------------------------------------
CLOUD_PATHS = [
    {
        "destination": "Oracle Database@Azure",
        "tag": "ANCHOR",
        "best_for": "Oracle-faithful workloads — Exadata, Autonomous, RAC, Data Guard — that need Azure as home with zero re-platform.",
        "preserved": "All Oracle-native features, DBA workflows, licensing model.",
        "azure_revenue": "OD@A consumption (Exadata ratecards) + downstream Fabric / Foundry / Copilot pull-through.",
        "color": ORACLE_RED,
        "emphasised": True,
    },
    {
        "destination": "Oracle DB on OCI",
        "tag": "PEER PATH",
        "best_for": "Multicloud customers, sovereign workloads, deep OCI investments.",
        "preserved": "OCI-native services + tooling; full second-cloud SLA.",
        "azure_revenue": "Fabric + Foundry + Copilot via shortcuts / GoldenGate.",
        "color": AZURE_BLUE,
        "emphasised": False,
    },
    {
        "destination": "Azure DB for PostgreSQL",
        "tag": "PEER PATH",
        "best_for": "Open-source-compatible workloads where re-platform is the right answer.",
        "preserved": "Application logic (with schema conversion); cloud-native ops.",
        "azure_revenue": "Azure PG consumption + downstream Fabric / Foundry / Copilot.",
        "color": IQ_YELLOW,
        "emphasised": False,
    },
    {
        "destination": "Stay on-prem (hybrid)",
        "tag": "PEER PATH",
        "best_for": "Stranded, regulated, or sovereign workloads that cannot leave the datacentre.",
        "preserved": "On-prem investment, data-residency posture, regulatory scope.",
        "azure_revenue": "Fabric capacity + GoldenGate + Foundry + Copilot via hybrid mirror.",
        "color": MUTED,
        "emphasised": False,
    },
]


def cloud_path_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE CLOUD-PATH MATRIX",
        size=11, color=ORACLE_RED, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Four legitimate destinations. OD@A is the anchor; the rest are honest peer paths.",
        size=24, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.5),
        "The honest framing — and the only credible way to position Azure as THE destination for Oracle workloads — is to give every workload a destination that actually fits.",
        size=12, color=WHITE_70,
    )

    # 4 columns. OD@A is wider (anchor emphasis); other 3 are narrower peers.
    col_top = Inches(2.2)
    col_h = Inches(4.6)
    total_w = Inches(12.1)
    # Allocate: 3.7 (OD@A) + 2.7 + 2.7 + 2.7 + 3 gaps of 0.1
    anchor_w = Inches(3.7)
    peer_w = Inches(2.7)
    gap = Inches(0.1)
    left = Inches(0.6)

    widths = [anchor_w, peer_w, peer_w, peer_w]
    for i, path in enumerate(CLOUD_PATHS):
        col_w = widths[i]
        color = path["color"]
        is_anchor = path["emphasised"]

        # Card — anchor is darker fill + heavier border
        card = add_card(slide, left, col_top, col_w, col_h, border=color, fill=NAVY_700,
                        fill_alpha=0.3 if is_anchor else 0.6)
        if is_anchor:
            card.line.width = Pt(2.25)

        # Tag pill at top
        tag_w = Inches(1.1) if not is_anchor else Inches(1.0)
        tag_left = left + Inches(0.18)
        add_pill(slide, tag_left, col_top + Inches(0.2), tag_w, Inches(0.32),
                 path["tag"], color)
        if is_anchor:
            add_pill(slide, tag_left + tag_w + Inches(0.1), col_top + Inches(0.2),
                     Inches(0.9), Inches(0.32), "OD@A", ORACLE_RED)

        # Destination name
        add_text(slide, left + Inches(0.18), col_top + Inches(0.6),
                 col_w - Inches(0.3), Inches(0.7),
                 path["destination"], size=15 if is_anchor else 13, bold=True, color=WHITE)

        # Body sections
        y = col_top + Inches(1.45)
        sections = [
            ("BEST FOR", path["best_for"], MUTED, WHITE_85),
            ("WHAT'S PRESERVED", path["preserved"], MUTED, WHITE_85),
            ("AZURE REVENUE LINE", path["azure_revenue"], MUTED, color),
        ]
        for label, body, label_color, body_color in sections:
            add_text(slide, left + Inches(0.18), y,
                     col_w - Inches(0.3), Inches(0.22),
                     label, size=8.5, color=label_color, font="Consolas")
            add_text(slide, left + Inches(0.18), y + Inches(0.24),
                     col_w - Inches(0.3), Inches(0.85),
                     body, size=10, color=body_color, bold=(label == "AZURE REVENUE LINE"))
            y += Inches(1.05)

        left += col_w + gap

    add_footer(slide=slide, prs=prs,
               text="Reference: /oda-frontier  ·  Cloud-path matrix · §5")


# ---------------------------------------------------------------------------
# Slide 5 — Technical architecture (5-layer summary)
# ---------------------------------------------------------------------------
ARCH_LAYERS = [
    ("VALUE LAYER · AGENTIC IQ STACK",
     "Fabric Mirror  ·  Fabric IQ  ·  Foundry IQ  ·  Work IQ  ·  Microsoft Copilot",
     IQ_TEAL),
    ("APP MODERNISATION",
     "AKS  ·  App Service  ·  Functions  ·  Entra ID  ·  Key Vault  ·  Application Insights",
     IQ_YELLOW),
    ("CLOUD LANDING ZONES",
     "Oracle Database@Azure (ANCHOR)   |   Oracle DB on OCI   |   Azure DB for PostgreSQL   |   On-prem (hybrid mirror)",
     ORACLE_RED),
    ("MIGRATION + UPGRADE",
     "Azure Migrate  ·  Oracle ZDM  ·  GoldenGate  ·  Azure DMS  ·  Compliant-version upgrade gates",
     AZURE_BLUE),
    ("CUSTOMER ESTATE",
     "Exadata  ·  RAC + Data Guard  ·  Database Appliance  ·  On-prem Fusion Apps",
     MUTED),
]


def architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "TECHNICAL ARCHITECTURE",
        size=11, color=AZURE_BLUE, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "End-to-end stack: estate at the bottom, OD@A in the middle, IQ stack on top.",
        size=24, bold=True, color=WHITE,
    )

    # Five stacked bands, bottom-up
    band_h = Inches(0.85)
    band_w = Inches(12.1)
    band_gap = Inches(0.08)
    bottom_y = Inches(6.0)

    for i, (label, body, color) in enumerate(ARCH_LAYERS):
        y = bottom_y - i * (band_h + band_gap)
        emphasised = i == 2  # cloud landing zones
        card = add_card(slide, Inches(0.6), y, band_w, band_h, border=color, fill=NAVY_700,
                        fill_alpha=0.4 if emphasised else 0.65)
        if emphasised:
            card.line.width = Pt(2.25)
        # Layer label
        add_text(slide, Inches(0.78), y + Inches(0.12), Inches(4.5), Inches(0.3),
                 label, size=10, color=color, font="Consolas", bold=True)
        # Body
        add_text(slide, Inches(0.78), y + Inches(0.42), band_w - Inches(0.4), Inches(0.4),
                 body, size=13 if emphasised else 11.5,
                 color=WHITE if emphasised else WHITE_85,
                 bold=emphasised)

    # Governance side note (lateral)
    gov_top = bottom_y + band_h + Inches(0.12)
    add_text(slide, Inches(0.6), gov_top, Inches(12.1), Inches(0.3),
             "Governance band runs laterally: Oracle Data Safe + Audit Vault  ↔  Microsoft Purview  —  lineage end-to-end.",
             size=11, color=MUTED, align=PP_ALIGN.CENTER, font="Consolas")

    add_footer(slide=slide, prs=prs,
               text="Reference: /oda-frontier  ·  Technical architecture · §6")


# ---------------------------------------------------------------------------
# Slide 6 — Consumption story (per-stage revenue table)
# ---------------------------------------------------------------------------
CONSUMPTION_ROWS = [
    ("01", "Discover", "Azure Migrate · Purview scan capacity", "Cloud Lift Services engagement", MUTED),
    ("02", "Plan", "Microsoft + partner architecture services", "Joint architecture council engagement", AZURE_BLUE),
    ("03", "Migrate", "Azure DMS · ASR · target compute on landing", "ZDM · GoldenGate · Cloud Lift", ORACLE_RED),
    ("04", "Modernise", "AKS · App Service · Functions · Entra · Key Vault · App Insights", "OD@A consumption — Exadata · Autonomous · DB Service", IQ_YELLOW),
    ("05", "Activate AI", "Fabric · OneLake · Foundry · Agent 365 · Copilot · Purview", "Fusion AI Agents · Audit Vault · Data Safe", IQ_TEAL),
]


def consumption_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE CONSUMPTION STORY",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Each stage of the journey lights up a distinct revenue line on both sides.",
        size=24, bold=True, color=WHITE,
    )

    # Header row
    header_top = Inches(1.7)
    col_x = [Inches(0.6), Inches(2.1), Inches(6.2), Inches(10.3)]
    col_w = [Inches(1.4), Inches(4.0), Inches(4.0), Inches(2.3)]

    headers = [("STAGE", MUTED), ("MICROSOFT REVENUE", AZURE_BLUE), ("ORACLE REVENUE", ORACLE_RED), ("STAGE ACCENT", MUTED)]
    for (label, color), x, w in zip(headers, col_x, col_w):
        add_text(slide, x, header_top, w, Inches(0.3), label,
                 size=10, color=color, font="Consolas", bold=True)

    # Header underline
    underline = slide.shapes.add_connector(1, Inches(0.6), header_top + Inches(0.36),
                                           Inches(12.7), header_top + Inches(0.36))
    underline.line.color.rgb = MUTED
    underline.line.width = Pt(0.5)
    underline.line.fill.transparency = 0.4

    # Rows
    row_top = Inches(2.2)
    row_h = Inches(0.85)
    for i, (num, stage, ms_rev, ora_rev, color) in enumerate(CONSUMPTION_ROWS):
        y = row_top + i * row_h
        # Subtle band on alternating rows
        if i % 2 == 0:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y,
                                          Inches(12.1), row_h - Inches(0.05))
            band.line.fill.background()
            band.fill.solid()
            band.fill.fore_color.rgb = NAVY_700
            band.fill.transparency = 0.55

        # Stage column
        add_text(slide, col_x[0], y + Inches(0.18), Inches(0.4), Inches(0.4),
                 num, size=15, bold=True, color=color, font="Consolas")
        add_text(slide, col_x[0] + Inches(0.45), y + Inches(0.2), col_w[0] - Inches(0.5), Inches(0.4),
                 stage, size=13, bold=True, color=WHITE)
        # Microsoft revenue
        add_text(slide, col_x[1], y + Inches(0.18), col_w[1], Inches(0.6),
                 ms_rev, size=11, color=AZURE_BLUE)
        # Oracle revenue
        add_text(slide, col_x[2], y + Inches(0.18), col_w[2], Inches(0.6),
                 ora_rev, size=11, color=ORACLE_RED)
        # Stage accent pill
        add_pill(slide, col_x[3], y + Inches(0.22), col_w[3], Inches(0.32),
                 f"ACT · {stage.upper()}", color)

    # Insight footer
    add_text(
        slide, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.45),
        "Stages 4 (modernise) and 5 (activate) are the compounding lines. OD@A consumption rises as more apps modernise on top; IQ-stack consumption rises as more workloads activate. Neither cannibalises the other.",
        size=11, color=WHITE_85, font="Calibri",
    )

    add_footer(slide=slide, prs=prs,
               text="Reference: /oda-frontier  ·  Consumption story · §7")


# ---------------------------------------------------------------------------
# Slide 7 — The exec ask
# ---------------------------------------------------------------------------
EXEC_ASKS = [
    (
        "Name an OD@A landing-zone owner per industry",
        "FSI, Manufacturing, Retail, Public Sector — one named technical sponsor each, accountable for the per-vertical reference architecture and the first three customer landings.",
        AZURE_BLUE,
    ),
    (
        "Fund the discovery accelerator",
        "Co-funded with Oracle Cloud Lift. Productised assessment that returns workload inventory, sensitivity catalogue, and per-workload destination recommendation in under 30 days.",
        IQ_YELLOW,
    ),
    (
        "Lock the cloud-path decision tree as the standard sales asset",
        "Make the four-path matrix (OD@A · OCI · Azure PG · hybrid) the canonical answer when the field is asked 'what does Microsoft recommend?' for an Oracle workload.",
        IQ_TEAL,
    ),
    (
        "Partner co-funded migration motion",
        "Accenture, Deloitte, EY, Capgemini packaged offers for ZDM-led migrations. Bundle with Cloud Lift credits and Universal Cloud Credits where they apply.",
        ORACLE_RED,
    ),
]


def ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE EXEC ASK",
        size=11, color=IQ_YELLOW, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "What lighting up OD@A as the anchor requires from Microsoft leadership.",
        size=24, bold=True, color=WHITE,
    )

    # 2x2 grid of asks
    grid_top = Inches(1.9)
    card_w = Inches(6.0)
    card_h = Inches(2.0)
    gap = Inches(0.2)
    left_start = Inches(0.6)

    for i, (title, body, color) in enumerate(EXEC_ASKS):
        col = i % 2
        row = i // 2
        cl = left_start + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)

        add_card(slide, cl, ct, card_w, card_h, border=color, fill=NAVY_700)
        # Number
        add_text(slide, cl + Inches(0.2), ct + Inches(0.2), Inches(0.6), Inches(0.4),
                 f"0{i + 1}", size=18, bold=True, color=color, font="Consolas")
        # Title
        add_text(slide, cl + Inches(0.85), ct + Inches(0.22), card_w - Inches(1.0), Inches(0.45),
                 title, size=14, bold=True, color=WHITE)
        # Body
        add_text(slide, cl + Inches(0.2), ct + Inches(0.85), card_w - Inches(0.4), Inches(1.1),
                 body, size=11, color=WHITE_85)

    # The bet — closing band
    bet_top = Inches(6.3)
    bet_card = add_card(slide, Inches(0.6), bet_top, Inches(12.1), Inches(0.85),
                        border=IQ_YELLOW, fill=NAVY_700, fill_alpha=0.4)
    bet_card.line.width = Pt(2)
    add_text(slide, Inches(0.85), bet_top + Inches(0.1), Inches(2), Inches(0.3),
             "THE BET", size=10, bold=True, color=IQ_YELLOW, font="Consolas")
    add_text(slide, Inches(0.85), bet_top + Inches(0.35), Inches(11.6), Inches(0.5),
             "Make OD@A the default destination Microsoft recommends for Oracle workloads, and the agentic IQ stack the default activation layer on top. Every other path remains available — OD@A is the anchor that compounds with Fabric, Foundry, Copilot, and Purview.",
             size=11.5, color=WHITE_85)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    anchor_slide(prs)
    journey_slide(prs)
    cloud_path_slide(prs)
    architecture_slide(prs)
    consumption_slide(prs)
    ask_slide(prs)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "decks"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "oda-frontier.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build()
