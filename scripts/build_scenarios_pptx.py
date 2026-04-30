"""Build a downloadable scenarios deck for the Enterprise Agent Alliance demo.

Output: docs/decks/scenarios.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Brand palette (matches the portal Tailwind config)
NAVY_900 = RGBColor(0x0B, 0x12, 0x20)
NAVY_700 = RGBColor(0x1A, 0x25, 0x40)
ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
IQ_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
IQ_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_85 = RGBColor(0xD9, 0xDD, 0xE5)
WHITE_70 = RGBColor(0xB0, 0xB7, 0xC4)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

TIER_COLORS = {
    "T0": MUTED,
    "T1": AZURE_BLUE,
    "T2": IQ_YELLOW,
    "T3": IQ_TEAL,
}


def add_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_900
    bg.shadow.inherit = False


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_pill(slide, left, top, width, height, text, color):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(0.75)
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.fill.transparency = 0.85
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Consolas"


def add_card(slide, left, top, width, height, *, border=MUTED, fill=NAVY_700, fill_alpha=0.6):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.line.color.rgb = border
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.fill.transparency = fill_alpha
    card.shadow.inherit = False
    return card


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # Eyebrow
    add_text(
        slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
        "THE ENTERPRISE AGENT ALLIANCE  ·  ORACLE × MICROSOFT",
        size=11, color=IQ_YELLOW, font="Consolas",
    )

    # Title
    add_text(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(2.2),
        "Same Oracle data.\nSame prompt.\nFour very different answers.",
        size=44, bold=True, color=WHITE,
    )

    # Subtitle
    add_text(
        slide, Inches(0.6), Inches(4.6), Inches(12), Inches(1.8),
        "Three real-world scenarios — Supply Chain, Workforce, Customer Health — and how Microsoft "
        "Copilot's answer evolves as Fabric IQ, Foundry IQ, and Work IQ are layered on top of "
        "Oracle data.",
        size=18, color=WHITE_85,
    )

    # Footer band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5),
        prs.slide_width, Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.45), Inches(12), Inches(0.4),
        "Oracle anchors data of record where it's deployed.  ·  Microsoft extends the worker surface.  ·  Together: agentic transformation.",
        size=10, color=MUTED, font="Consolas",
    )


def framing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
        "HOW TO READ THESE SCENARIOS",
        size=11, color=MUTED, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.7),
        "Each scenario asks the same question across four tiers.",
        size=26, bold=True, color=WHITE,
    )

    rows = [
        ("T0", "Baseline", "Lookup agent",
         "Knows where Oracle data lives. Describes tables. Cannot interpret.",
         "$0 incremental decision value. Foundation only."),
        ("T1", "+ Fabric IQ", "Insight agent",
         "Joins the semantic model. Quantifies, aggregates, trends.",
         "Defensible numbers without SQL. Reporting goes from days to minutes."),
        ("T2", "+ Foundry IQ", "Decision agent",
         "Reasons across the data. Ranks moves with projected $-impact and grounded citations.",
         "Replaces an analyst week with a conversation."),
        ("T3", "+ Work IQ", "Execution agent",
         "Personalised via Microsoft Graph. Drafts the email, the Teams message, the meeting agenda.",
         "The decision is in the calendar before the meeting starts."),
    ]

    top = Inches(2.0)
    row_h = Inches(1.05)
    left = Inches(0.6)
    width = Inches(12.1)

    for i, (tier, label, archetype, delivers, benefit) in enumerate(rows):
        row_top = top + i * row_h
        color = TIER_COLORS[tier]
        add_card(slide, left, row_top, width, row_h - Inches(0.1), border=color)
        # Tier badge
        add_text(slide, left + Inches(0.2), row_top + Inches(0.18), Inches(0.7), Inches(0.3),
                 tier, size=12, bold=True, color=color, font="Consolas")
        add_text(slide, left + Inches(0.85), row_top + Inches(0.18), Inches(2.2), Inches(0.3),
                 label, size=13, bold=True, color=WHITE)
        add_text(slide, left + Inches(0.85), row_top + Inches(0.5), Inches(2.2), Inches(0.3),
                 archetype, size=10, color=color, font="Consolas")
        # Delivers
        add_text(slide, left + Inches(3.2), row_top + Inches(0.18), Inches(5.4), Inches(0.7),
                 delivers, size=11, color=WHITE_85)
        # Benefit
        add_text(slide, left + Inches(8.7), row_top + Inches(0.18), Inches(3.3), Inches(0.7),
                 benefit, size=11, color=color)

    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.6), Inches(12), Inches(0.4),
        "By Tier 3, the agent isn't reporting on the business — it's running parts of it.",
        size=12, color=MUTED, font="Calibri",
    )


def scenario_slide(prs, *, eyebrow, title, prompt, accent, tiers):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # Eyebrow
    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        eyebrow, size=10, color=accent, font="Consolas",
    )
    # Title
    add_text(
        slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.6),
        title, size=28, bold=True, color=WHITE,
    )
    # Prompt
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12), Inches(0.6),
        f"“{prompt}”",
        size=14, color=WHITE_85,
    )

    # 4 tier columns
    col_top = Inches(2.4)
    col_h = Inches(4.6)
    col_w = Inches(2.95)
    gap = Inches(0.15)
    left_start = Inches(0.6)

    for i, (tier_id, tier_label, archetype, delivers, benefit) in enumerate(tiers):
        col_left = left_start + i * (col_w + gap)
        color = TIER_COLORS[tier_id]
        add_card(slide, col_left, col_top, col_w, col_h, border=color, fill=NAVY_700)

        # Tier id chip
        add_text(slide, col_left + Inches(0.18), col_top + Inches(0.18), Inches(2.6), Inches(0.3),
                 f"{tier_id}  ·  {tier_label}", size=10, bold=True, color=color, font="Consolas")
        # Archetype
        add_text(slide, col_left + Inches(0.18), col_top + Inches(0.55), Inches(2.6), Inches(0.4),
                 archetype, size=14, bold=True, color=WHITE)

        # Divider line
        line = slide.shapes.add_connector(1, col_left + Inches(0.18), col_top + Inches(1.05),
                                          col_left + col_w - Inches(0.18), col_top + Inches(1.05))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6

        # Delivers label + body
        add_text(slide, col_left + Inches(0.18), col_top + Inches(1.15), Inches(2.6), Inches(0.25),
                 "DELIVERS", size=9, color=MUTED, font="Consolas")
        add_text(slide, col_left + Inches(0.18), col_top + Inches(1.4), Inches(2.6), Inches(1.5),
                 delivers, size=11, color=WHITE_85)

        # Benefit label + body
        add_text(slide, col_left + Inches(0.18), col_top + Inches(3.0), Inches(2.6), Inches(0.25),
                 "BUSINESS BENEFIT", size=9, color=MUTED, font="Consolas")
        add_text(slide, col_left + Inches(0.18), col_top + Inches(3.25), Inches(2.6), Inches(1.3),
                 benefit, size=11, color=color, bold=True)


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
        "THE COMPOUNDING EFFECT",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(1.1), Inches(12), Inches(1.5),
        "Most agentic projects ship at Tier 1 — they report, but never act.",
        size=28, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(2.6), Inches(12), Inches(2.0),
        "Microsoft and Oracle together ship the full progression.\n\n"
        "Oracle anchors the data of record where it's deployed.\n"
        "Microsoft extends the worker surface where the work actually happens.\n"
        "Together: the agent isn't reporting on the business — it's running parts of it.",
        size=16, color=WHITE_85,
    )

    # Three pills
    base_top = Inches(5.6)
    pill_w = Inches(3.9)
    add_pill(slide, Inches(0.6), base_top, pill_w, Inches(0.5),
             "ORACLE  ·  Data of record", ORACLE_RED)
    add_pill(slide, Inches(4.7), base_top, pill_w, Inches(0.5),
             "MICROSOFT  ·  System of work", AZURE_BLUE)
    add_pill(slide, Inches(8.8), base_top, pill_w, Inches(0.5),
             "AGENTIC TRANSFORMATION", IQ_TEAL)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    framing_slide(prs)

    scenario_slide(
        prs,
        eyebrow="ORACLE FUSION SCM  ·  PROCUREMENT",
        title="Scenario 1 — Supply Chain Risk",
        prompt="What's our Q3 supply chain risk exposure in EMEA, and which suppliers should we de-risk first?",
        accent=AZURE_BLUE,
        tiers=[
            ("T0", "Baseline", "Lookup agent",
             "Lists the Oracle SCM tables that hold suppliers, POs, and delivery records.",
             "Foundation in place — Oracle data is mirrored and ready, but answers stay generic."),
            ("T1", "+ Fabric IQ", "Insight agent",
             "$47.2M EMEA exposure quantified across 84 suppliers; OTD has fallen to 87%.",
             "Risk visibility moves from anecdotes to a defensible number, instantly."),
            ("T2", "+ Foundry IQ", "Decision agent",
             "5 prioritised de-risk plays with $18.6M mitigatable inside Q3. First move: dual-source 25% of Adriatica volume.",
             "Action plan with projected savings — no analyst week required."),
            ("T3", "+ Work IQ", "Execution agent",
             "Email to Maria drafted; Teams ping to Tomás drafted; tomorrow's 10am EMEA ops agenda pre-filled.",
             "Decisions executed at the speed of a chat; supplier risk doesn't wait a week."),
        ],
    )

    scenario_slide(
        prs,
        eyebrow="ORACLE FUSION HCM  ·  RECRUITING",
        title="Scenario 2 — Workforce Planning",
        prompt="Summarise open requisitions on my team and flag any at risk of missing the fiscal-year close.",
        accent=IQ_YELLOW,
        tiers=[
            ("T0", "Baseline", "Lookup agent",
             "Names the HCM tables for requisitions, pipeline, and offers.",
             "Knows the data exists — but can't yet resolve which reqs are 'on your team.'"),
            ("T1", "+ Fabric IQ", "Insight agent",
             "14 open reqs · 5 at FY-close risk · pipeline coverage has fallen to 3.4×.",
             "FY hiring health visible on one screen, scoped to your org rollup."),
            ("T2", "+ Foundry IQ", "Decision agent",
             "Per-req recommendation: re-scope, push to offer, or open to internal mobility.",
             "Hiring plan that actually closes by FY end — projected close rate +18 pts."),
            ("T3", "+ Work IQ", "Execution agent",
             "Friday review with Diego pre-filled; 1:1 agenda for Aisha (EM candidate) drafted.",
             "Manager time goes to deciding, not to prepping pre-reads."),
        ],
    )

    scenario_slide(
        prs,
        eyebrow="ORACLE FUSION CX  ·  SERVICE",
        title="Scenario 3 — Customer Health",
        prompt="Which accounts in my book are showing churn signals, and what's the recommended next play for each?",
        accent=IQ_TEAL,
        tiers=[
            ("T0", "Baseline", "Lookup agent",
             "Points at the CX tables for accounts, health, and service requests.",
             "Foundation only — can't yet score churn or scope to your book of business."),
            ("T1", "+ Fabric IQ", "Insight agent",
             "32 accounts · 4 at risk · $6.0M ARR exposed.",
             "Churn risk surfaces early, with the dollars attached."),
            ("T2", "+ Foundry IQ", "Decision agent",
             "Right play per account; $4.6M of the $6.0M is recoverable.",
             "CSMs lead with the highest-yield save play, not a generic check-in."),
            ("T3", "+ Work IQ", "Execution agent",
             "Northwind QBR pre-loaded with the value review; Contoso CXO call prep with Liam ready.",
             "The save motion is in the calendar before the customer calls you."),
        ],
    )

    closing_slide(prs)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "decks"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "scenarios.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build()
