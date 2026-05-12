"""Build the Frontier Transformation exec deck — the stakeholder briefing
that mirrors the /frontier-transformation portal page.

Output: docs/decks/frontier-transformation.pptx

Seven slides:
  1. Title — "Three coordinated workstreams. One outcome: the frontier firm
     running on Azure."
  2. The three workstreams (overview)
  3. THE COMPLETE ARCHITECTURE — single show-piece slide with the entire
     stack: goal -> three workstream columns -> interop lanes -> outcome.
  4. Joint agent ecosystem — MCP, A2A, Fabric Mirror interop between
     Microsoft and Oracle agents.
  5. Workload Virtual Instance + AIOps stack — WVI as the unit of AIOps,
     with the Common Agent Platform and its agents.
  6. The unified outcome — four stakeholder lenses on the frontier firm.
  7. The ask — four commits.

Visual conventions, brand palette, and helper functions mirror
scripts/build_scenarios_pptx.py, scripts/build_oda_frontier_pptx.py,
scripts/build_exec_summary_pptx.py so all four decks read as a family.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette (matches the portal + sibling decks)
# ---------------------------------------------------------------------------
NAVY_900 = RGBColor(0x0B, 0x12, 0x20)
NAVY_700 = RGBColor(0x1A, 0x25, 0x40)
ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
IQ_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
IQ_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_85 = RGBColor(0xD9, 0xDD, 0xE5)
WHITE_70 = RGBColor(0xB0, 0xB7, 0xC4)
MUTED = RGBColor(0x94, 0xA3, 0xB8)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_900
    bg.shadow.inherit = False


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_pill(slide, left, top, width, height, text, color, *, size=9, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(0.75)
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.fill.transparency = 0.85
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Consolas"


def add_card(slide, left, top, width, height, *, border=MUTED, fill=NAVY_700, fill_alpha=0.6):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.line.color.rgb = border
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.fill.transparency = fill_alpha
    card.shadow.inherit = False
    return card


def add_footer(slide, prs, text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4),
        prs.slide_width, Inches(0.4),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.36), Inches(12), Inches(0.32),
        text, size=9, color=MUTED, font="Consolas",
    )


def add_solid(slide, left, top, width, height, color, *, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.18
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
        "FRONTIER TRANSFORMATION  ·  BIZ OPS + INFRA OPS  ·  ORACLE × MICROSOFT",
        size=11, color=IQ_TEAL, font="Consolas",
    )

    add_text(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(2.8),
        "Three coordinated workstreams.\nOne outcome:\nthe frontier firm running on Azure.",
        size=38, bold=True, color=WHITE,
    )

    add_text(
        slide, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8),
        "OD@A anchors the data estate. The Microsoft Agent Framework modernises business ops. "
        "The Workload Virtual Instance is the new unit of AIOps for infra ops. "
        "Microsoft and Oracle agents interoperate over MCP, A2A, and Fabric Mirror across all three.",
        size=17, color=WHITE_85,
    )

    # Footer band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5),
        prs.slide_width, Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.45), Inches(12), Inches(0.4),
        "Make Azure the best place to run Oracle workloads — across data, biz ops, and infra ops.",
        size=10, color=MUTED, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 2 — The three workstreams (overview)
# ---------------------------------------------------------------------------
WORKSTREAMS_DECK = [
    {
        "num": "01",
        "label": "Data Estate",
        "agents": "JOINT",
        "title": "Modernise the data estate on Azure",
        "body": "Every Oracle workload to its right destination — OD@A as the anchor, with peer paths for OCI, Azure DB for PostgreSQL, and hybrid mirror.",
        "outcome": "Data of record anchored on Azure. The agentic stack grounds wherever it lives.",
        "color": ORACLE_RED,
    },
    {
        "num": "02",
        "label": "Biz Ops",
        "agents": "MICROSOFT",
        "title": "Modernise the apps + ship agents",
        "body": "App stack to cloud standards. Agents on the Microsoft Agent Framework — rooted in the IQ layer, A365, and Purview — interoperating with Oracle agents.",
        "outcome": "Worker-facing agents that draft, decide, and execute across the modernised app surface.",
        "color": AZURE_BLUE,
    },
    {
        "num": "03",
        "label": "Infra Ops",
        "agents": "JOINT",
        "title": "Workload Virtual Instance + AIOps",
        "body": "WVI as the new unit of AIOps. Common Agent Platform shipping Config Drift, Troubleshooting & RCA agents on a Quality / Resiliency / Security foundation.",
        "outcome": "Drift caught before incidents. RCA at machine speed. Workload health green by default.",
        "color": IQ_TEAL,
    },
]


def workstreams_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE FRAME · THREE WORKSTREAMS",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Three coordinated workstreams. Each partial alone. Together they ship the frontier firm.",
        size=22, bold=True, color=WHITE,
    )

    # Three workstream columns
    col_top = Inches(1.85)
    col_h = Inches(4.9)
    col_w = Inches(4.1)
    gap = Inches(0.15)
    left = Inches(0.6)

    for i, w in enumerate(WORKSTREAMS_DECK):
        cl = left + i * (col_w + gap)
        color = w["color"]
        add_card(slide, cl, col_top, col_w, col_h, border=color, fill=NAVY_700)

        # Number + label + agents pill
        add_text(slide, cl + Inches(0.25), col_top + Inches(0.2),
                 Inches(0.9), Inches(0.45),
                 w["num"], size=20, bold=True, color=color, font="Consolas")
        add_pill(slide, cl + col_w - Inches(1.25), col_top + Inches(0.27),
                 Inches(1.0), Inches(0.3), w["agents"], color, size=9)

        add_text(slide, cl + Inches(0.25), col_top + Inches(0.75),
                 col_w - Inches(0.5), Inches(0.45),
                 w["label"].upper(), size=11, bold=True, color=color, font="Consolas")
        add_text(slide, cl + Inches(0.25), col_top + Inches(1.05),
                 col_w - Inches(0.5), Inches(0.6),
                 w["title"], size=17, bold=True, color=WHITE)

        # Body
        add_text(slide, cl + Inches(0.25), col_top + Inches(1.95),
                 col_w - Inches(0.5), Inches(1.65),
                 w["body"], size=12, color=WHITE_85)

        # Outcome
        line = slide.shapes.add_connector(1, cl + Inches(0.25), col_top + Inches(3.55),
                                          cl + col_w - Inches(0.25), col_top + Inches(3.55))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        add_text(slide, cl + Inches(0.25), col_top + Inches(3.6),
                 col_w - Inches(0.5), Inches(0.25),
                 "OUTCOME", size=9, color=MUTED, font="Consolas")
        add_text(slide, cl + Inches(0.25), col_top + Inches(3.85),
                 col_w - Inches(0.5), Inches(1.0),
                 w["outcome"], size=11.5, color=color, bold=True)

    # Bottom unified outcome strip
    band = add_solid(slide, Inches(0.6), Inches(6.9), Inches(12.1), Inches(0.45),
                     IQ_TEAL, line=IQ_TEAL)
    band.fill.transparency = 0.65
    add_text(slide, Inches(0.85), Inches(6.96), Inches(11.6), Inches(0.35),
             "UNIFIED OUTCOME  ·  Frontier firm — agent-augmented, data-grounded, cross-stack, governed by default.",
             size=11, color=IQ_TEAL, font="Consolas", bold=True)


# ---------------------------------------------------------------------------
# Slide 3 — THE COMPLETE ARCHITECTURE (the show-piece)
# ---------------------------------------------------------------------------
def complete_arch_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    # Eyebrow
    add_text(
        slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.3),
        "THE COMPLETE ARCHITECTURE  ·  ALL THREE WORKSTREAMS IN ONE FRAME",
        size=11, color=IQ_TEAL, font="Consolas",
    )

    # Goal pill at the very top
    goal_card = add_solid(slide, Inches(0.4), Inches(0.55), Inches(12.5), Inches(0.4),
                          IQ_YELLOW, line=IQ_YELLOW)
    goal_card.fill.transparency = 0.6
    add_text(slide, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.3),
             "GOAL  ·  MAKE AZURE THE BEST PLACE TO RUN ORACLE WORKLOADS",
             size=11.5, bold=True, color=IQ_YELLOW, font="Consolas",
             align=PP_ALIGN.CENTER)

    # Three workstream columns
    col_top = Inches(1.1)
    col_h = Inches(5.2)
    col_w = Inches(4.05)
    gap = Inches(0.15)
    left_start = Inches(0.4)

    # ----- Column 1: Data Estate -----
    c1_x = left_start
    add_card(slide, c1_x, col_top, col_w, col_h, border=ORACLE_RED, fill=NAVY_700)
    add_text(slide, c1_x + Inches(0.2), col_top + Inches(0.15),
             Inches(0.6), Inches(0.3),
             "01", size=14, bold=True, color=ORACLE_RED, font="Consolas")
    add_text(slide, c1_x + Inches(0.65), col_top + Inches(0.15),
             col_w - Inches(0.8), Inches(0.3),
             "DATA ESTATE  ·  WORKSTREAM 1", size=10, bold=True, color=ORACLE_RED, font="Consolas")
    add_text(slide, c1_x + Inches(0.2), col_top + Inches(0.45),
             col_w - Inches(0.4), Inches(0.35),
             "Modernise on OD@A", size=15, bold=True, color=WHITE)

    # Stack inside column 1
    y = col_top + Inches(0.95)
    # Estate row
    bx = add_solid(slide, c1_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.4),
                   ORACLE_RED, line=ORACLE_RED)
    bx.fill.transparency = 0.78
    add_text(slide, c1_x + Inches(0.2), y + Inches(0.05),
             col_w - Inches(0.4), Inches(0.3),
             "On-prem estate · Exadata · RAC · DBA · Fusion Apps",
             size=10, color=WHITE_85, align=PP_ALIGN.CENTER, font="Consolas")
    y += Inches(0.5)
    # Up arrow
    add_text(slide, c1_x, y - Inches(0.1), col_w, Inches(0.2),
             "↑", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # Migration row
    bx = add_solid(slide, c1_x + Inches(0.2), y + Inches(0.05), col_w - Inches(0.4), Inches(0.4),
                   AZURE_BLUE, line=AZURE_BLUE)
    bx.fill.transparency = 0.78
    add_text(slide, c1_x + Inches(0.2), y + Inches(0.1),
             col_w - Inches(0.4), Inches(0.3),
             "Azure Migrate · ZDM · GoldenGate · Azure DMS",
             size=10, color=WHITE_85, align=PP_ALIGN.CENTER, font="Consolas")
    y += Inches(0.6)
    add_text(slide, c1_x, y - Inches(0.1), col_w, Inches(0.2),
             "↑", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # OD@A anchor (emphasised)
    anchor = add_solid(slide, c1_x + Inches(0.2), y + Inches(0.05), col_w - Inches(0.4), Inches(0.65),
                       ORACLE_RED, line=ORACLE_RED)
    anchor.line.width = Pt(2.25)
    anchor.fill.transparency = 0.55
    add_text(slide, c1_x + Inches(0.2), y + Inches(0.1),
             col_w - Inches(0.4), Inches(0.3),
             "Oracle Database@Azure",
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, c1_x + Inches(0.2), y + Inches(0.4),
             col_w - Inches(0.4), Inches(0.3),
             "ANCHOR · Exadata · Autonomous · DB Service",
             size=9.5, color=WHITE_85, align=PP_ALIGN.CENTER, font="Consolas")
    y += Inches(0.85)
    # Peer paths row (3 small pills)
    peer_w = (col_w - Inches(0.6)) / 3
    peers = [("OCI", AZURE_BLUE), ("Azure PG", IQ_YELLOW), ("Hybrid", MUTED)]
    for j, (lbl, col) in enumerate(peers):
        px = c1_x + Inches(0.2) + j * (peer_w + Inches(0.1))
        add_pill(slide, px, y + Inches(0.05), peer_w, Inches(0.3), lbl, col, size=8.5)
    y += Inches(0.5)
    add_text(slide, c1_x + Inches(0.2), y, col_w - Inches(0.4), Inches(0.25),
             "PEER PATHS", size=8.5, color=MUTED, font="Consolas",
             align=PP_ALIGN.CENTER)

    # ----- Column 2: Biz Ops -----
    c2_x = left_start + col_w + gap
    add_card(slide, c2_x, col_top, col_w, col_h, border=AZURE_BLUE, fill=NAVY_700)
    add_text(slide, c2_x + Inches(0.2), col_top + Inches(0.15),
             Inches(0.6), Inches(0.3),
             "02", size=14, bold=True, color=AZURE_BLUE, font="Consolas")
    add_text(slide, c2_x + Inches(0.65), col_top + Inches(0.15),
             col_w - Inches(0.8), Inches(0.3),
             "BIZ OPS  ·  WORKSTREAM 2", size=10, bold=True, color=AZURE_BLUE, font="Consolas")
    add_text(slide, c2_x + Inches(0.2), col_top + Inches(0.45),
             col_w - Inches(0.4), Inches(0.35),
             "Apps + Microsoft Agent Framework", size=15, bold=True, color=WHITE)

    biz_layers = [
        ("App stack · cloud-native",   "AKS · App Service · Functions · Entra · Key Vault",         IQ_YELLOW),
        ("Microsoft Agent Framework",  "Foundry · Copilot Studio · evals + tracing",                AZURE_BLUE),
        ("IQ layer",                   "Fabric IQ · Foundry IQ · Work IQ",                          IQ_TEAL),
        ("Microsoft Agent 365",        "Copilot · Outlook · Teams · control plane",                 AZURE_BLUE),
        ("Governance",                 "Microsoft Purview ↔ Oracle Data Safe",                      ORACLE_RED),
    ]
    layer_y = col_top + Inches(0.95)
    layer_h = Inches(0.78)
    for j, (lbl, sub, col) in enumerate(biz_layers):
        ly = layer_y + j * layer_h
        lc = add_solid(slide, c2_x + Inches(0.2), ly, col_w - Inches(0.4), Inches(0.68),
                       col, line=col)
        lc.fill.transparency = 0.78
        add_text(slide, c2_x + Inches(0.3), ly + Inches(0.05),
                 col_w - Inches(0.5), Inches(0.3),
                 lbl, size=11, bold=True, color=WHITE)
        add_text(slide, c2_x + Inches(0.3), ly + Inches(0.32),
                 col_w - Inches(0.5), Inches(0.3),
                 sub, size=9, color=WHITE_85, font="Consolas")

    # ----- Column 3: Infra Ops (WVI + AIOps) -----
    c3_x = left_start + 2 * (col_w + gap)
    add_card(slide, c3_x, col_top, col_w, col_h, border=IQ_TEAL, fill=NAVY_700)
    add_text(slide, c3_x + Inches(0.2), col_top + Inches(0.15),
             Inches(0.6), Inches(0.3),
             "03", size=14, bold=True, color=IQ_TEAL, font="Consolas")
    add_text(slide, c3_x + Inches(0.65), col_top + Inches(0.15),
             col_w - Inches(0.8), Inches(0.3),
             "INFRA OPS  ·  WORKSTREAM 3", size=10, bold=True, color=IQ_TEAL, font="Consolas")
    add_text(slide, c3_x + Inches(0.2), col_top + Inches(0.45),
             col_w - Inches(0.4), Inches(0.35),
             "WVI + AIOps stack", size=15, bold=True, color=WHITE)

    # WVI dashed envelope
    wvi_top = col_top + Inches(0.95)
    wvi_h = Inches(4.05)
    wvi = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  c3_x + Inches(0.15), wvi_top,
                                  col_w - Inches(0.3), wvi_h)
    wvi.adjustments[0] = 0.06
    wvi.line.color.rgb = IQ_TEAL
    wvi.line.width = Pt(1.5)
    wvi.line.dash_style = 7  # dashed
    wvi.fill.solid()
    wvi.fill.fore_color.rgb = IQ_TEAL
    wvi.fill.transparency = 0.92

    # WVI label tag
    add_pill(slide, c3_x + Inches(0.3), wvi_top + Inches(0.05),
             Inches(1.4), Inches(0.28), "WVI", IQ_TEAL, size=9)

    # Inside WVI: AIOps stack (compact)
    sx = c3_x + Inches(0.3)
    sw = col_w - Inches(0.6)
    sy = wvi_top + Inches(0.45)

    # Top band — issue prevention
    b = add_solid(slide, sx, sy, sw, Inches(0.36), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.6
    add_text(slide, sx, sy + Inches(0.05), sw, Inches(0.3),
             "Proactive + reactive issue prevention", size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.46)

    # Two agents side by side
    half = (sw - Inches(0.1)) / 2
    cd = add_solid(slide, sx, sy, half, Inches(0.48), AZURE_BLUE, line=AZURE_BLUE)
    cd.fill.transparency = 0.78
    add_text(slide, sx, sy + Inches(0.04), half, Inches(0.22),
             "Config Drift Agent", size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sx, sy + Inches(0.25), half, Inches(0.22),
             "PROACTIVE", size=8, color=AZURE_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    rc = add_solid(slide, sx + half + Inches(0.1), sy, half, Inches(0.48), ORACLE_RED, line=ORACLE_RED)
    rc.fill.transparency = 0.78
    add_text(slide, sx + half + Inches(0.1), sy + Inches(0.04), half, Inches(0.22),
             "Troubleshooting + RCA", size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sx + half + Inches(0.1), sy + Inches(0.25), half, Inches(0.22),
             "REACTIVE", size=8, color=ORACLE_RED, align=PP_ALIGN.CENTER, font="Consolas")
    sy += Inches(0.58)

    # Common Agent Platform band
    b = add_solid(slide, sx, sy, sw, Inches(0.36), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.55
    add_text(slide, sx, sy + Inches(0.05), sw, Inches(0.3),
             "Common Agent Platform", size=10.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.46)

    # Three pillars
    pillars = ["Configuration Validation", "Workload Health", "Contextual Support"]
    pillar_w = (sw - Inches(0.2)) / 3
    for j, p in enumerate(pillars):
        px = sx + j * (pillar_w + Inches(0.1))
        bx = add_solid(slide, px, sy, pillar_w, Inches(0.42),
                       WHITE, line=MUTED)
        bx.fill.transparency = 0.88
        add_text(slide, px, sy + Inches(0.05), pillar_w, Inches(0.18),
                 p, size=8, color=AZURE_BLUE, align=PP_ALIGN.CENTER, bold=True)
    sy += Inches(0.5)

    # Proactive | Reactive split
    pa = add_solid(slide, sx, sy, half, Inches(0.3), AZURE_BLUE, line=AZURE_BLUE)
    pa.fill.transparency = 0.55
    add_text(slide, sx, sy + Inches(0.04), half, Inches(0.22),
             "Proactive", size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    re = add_solid(slide, sx + half + Inches(0.1), sy, half, Inches(0.3),
                   ORACLE_RED, line=ORACLE_RED)
    re.fill.transparency = 0.55
    add_text(slide, sx + half + Inches(0.1), sy + Inches(0.04), half, Inches(0.22),
             "Reactive", size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.4)

    # Foundation band
    b = add_solid(slide, sx, sy, sw, Inches(0.36), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.55
    add_text(slide, sx, sy + Inches(0.05), sw, Inches(0.3),
             "Quality · Resiliency · Security",
             size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ----- Interop lanes spanning all three columns -----
    interop_top = col_top + col_h + Inches(0.18)
    interop_band = add_solid(slide, Inches(0.4), interop_top, Inches(12.5), Inches(0.5),
                             NAVY_700, line=MUTED)
    interop_band.fill.transparency = 0.4
    add_text(slide, Inches(0.5), interop_top + Inches(0.06), Inches(3), Inches(0.18),
             "AGENT INTEROP", size=9, color=MUTED, font="Consolas", bold=True)
    # Three lane pills inline
    lane_y = interop_top + Inches(0.05)
    add_pill(slide, Inches(3.4), lane_y, Inches(2.5), Inches(0.35),
             "MCP · MS → Oracle", IQ_TEAL, size=9)
    add_pill(slide, Inches(6.1), lane_y, Inches(2.7), Inches(0.35),
             "Fabric Mirror · Oracle → MS", AZURE_BLUE, size=9)
    add_pill(slide, Inches(9.0), lane_y, Inches(2.5), Inches(0.35),
             "A2A · Bidirectional", ORACLE_RED, size=9)
    add_text(slide, Inches(0.5), interop_top + Inches(0.28), Inches(3), Inches(0.18),
             "BETWEEN STACKS", size=8.5, color=MUTED, font="Consolas")

    # Outcome band at the very bottom
    outcome_top = interop_top + Inches(0.6)
    outcome = add_solid(slide, Inches(0.4), outcome_top, Inches(12.5), Inches(0.42),
                        IQ_TEAL, line=IQ_TEAL)
    outcome.fill.transparency = 0.55
    add_text(slide, Inches(0.4), outcome_top + Inches(0.05),
             Inches(12.5), Inches(0.32),
             "FRONTIER FIRM  ·  agent-augmented work  ·  data-grounded decisions  ·  cross-stack interop  ·  governed by default",
             size=11, bold=True, color=WHITE, font="Consolas", align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 4 — Joint agent ecosystem
# ---------------------------------------------------------------------------
INTEROP_LANES = [
    {
        "label": "MCP loop back to Oracle",
        "protocol": "Model Context Protocol",
        "direction": "MS → Oracle",
        "color": IQ_TEAL,
        "body": "Microsoft-side agent (Copilot / Foundry / A365) hands off to an Oracle Fusion AI Agent to execute a write-back, transaction, or policy-bound flow inside the Oracle estate.",
        "example": "Foundry recommends 'dual-source 25% of Adriatica volume' → MCP handoff → Fusion AI Agent executes the PO change.",
    },
    {
        "label": "Fabric Mirror (data plane)",
        "protocol": "Native Fabric Mirroring + GoldenGate",
        "direction": "Oracle → MS",
        "color": AZURE_BLUE,
        "body": "Oracle data of record streams into OneLake in near-real-time. Microsoft-side agents ground on the Fabric semantic model without ever copying tables.",
        "example": "Direct Lake query against mirrored Fusion SCM tables — sub-second-fresh, fully governed.",
    },
    {
        "label": "Agent-to-Agent (A2A) handoff",
        "protocol": "A2A · OpenLineage · OAuth 2.0",
        "direction": "Bidirectional",
        "color": ORACLE_RED,
        "body": "Microsoft agents and Oracle Fusion AI Agents call each other as peers. Cross-domain workflows in a single chat thread.",
        "example": "Work IQ + Foundry + Fusion HCM agents collaborate to fill an open req and update the requisition.",
    },
]


def interop_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "JOINT AGENT ECOSYSTEM",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Microsoft and Oracle agents call each other as peers.",
        size=24, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.4),
        "Three interop lanes: MCP for tool handoffs, Fabric Mirror for the data plane, A2A for true peer-to-peer agent orchestration.",
        size=12, color=WHITE_70,
    )

    # Three lane cards
    card_top = Inches(2.1)
    card_h = Inches(4.6)
    card_w = Inches(4.1)
    gap = Inches(0.15)
    left = Inches(0.6)

    for i, lane in enumerate(INTEROP_LANES):
        cl = left + i * (card_w + gap)
        color = lane["color"]
        add_card(slide, cl, card_top, card_w, card_h, border=color, fill=NAVY_700)
        # Direction pill at the top
        add_pill(slide, cl + Inches(0.2), card_top + Inches(0.2),
                 Inches(1.6), Inches(0.32), lane["direction"], color, size=9)
        # Protocol mono on the right
        add_text(slide, cl + Inches(1.9), card_top + Inches(0.22),
                 card_w - Inches(2.1), Inches(0.3),
                 lane["protocol"].upper(), size=8.5, color=MUTED, font="Consolas",
                 align=PP_ALIGN.RIGHT)
        # Title
        add_text(slide, cl + Inches(0.25), card_top + Inches(0.7),
                 card_w - Inches(0.5), Inches(0.6),
                 lane["label"], size=16, bold=True, color=WHITE)
        # Body
        add_text(slide, cl + Inches(0.25), card_top + Inches(1.5),
                 card_w - Inches(0.5), Inches(1.7),
                 lane["body"], size=12, color=WHITE_85)
        # Divider
        line = slide.shapes.add_connector(1, cl + Inches(0.25), card_top + Inches(3.2),
                                          cl + card_w - Inches(0.25), card_top + Inches(3.2))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        # Example
        add_text(slide, cl + Inches(0.25), card_top + Inches(3.3),
                 card_w - Inches(0.5), Inches(0.25),
                 "EXAMPLE", size=9, color=MUTED, font="Consolas")
        add_text(slide, cl + Inches(0.25), card_top + Inches(3.55),
                 card_w - Inches(0.5), Inches(1.0),
                 lane["example"], size=11, color=color, bold=True)


# ---------------------------------------------------------------------------
# Slide 5 — WVI + AIOps stack
# ---------------------------------------------------------------------------
WVI_PROPS = [
    ("Spans resource groups",
     "Logical envelope around everything one workload needs — across resource groups + subscriptions."),
    ("Unit of AIOps",
     "Agents target the WVI. Health, drift, RCA, remediation reasoned at the workload boundary."),
    ("Policy + governance anchor",
     "Compliance baselines, sovereign overlays, FedRAMP / IL5 profiles attach to the WVI."),
    ("Ownership clarity",
     "One workload, one owner, one SLA. No ambiguity when an agent flags an incident."),
]


def wvi_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "WORKLOAD VIRTUAL INSTANCE  ·  THE UNIT OF AIOPS",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "The Workload Virtual Instance becomes the new unit of AIOps.",
        size=24, bold=True, color=WHITE,
    )

    # Left: WVI properties (4 small cards in 2x2)
    grid_top = Inches(1.85)
    card_w = Inches(2.85)
    card_h = Inches(2.0)
    gap = Inches(0.15)
    left = Inches(0.6)

    for i, (lbl, body) in enumerate(WVI_PROPS):
        col = i % 2
        row = i // 2
        cl = left + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)
        add_card(slide, cl, ct, card_w, card_h, border=IQ_TEAL, fill=NAVY_700)
        add_text(slide, cl + Inches(0.2), ct + Inches(0.2),
                 card_w - Inches(0.4), Inches(0.35),
                 lbl.upper(), size=9.5, bold=True, color=IQ_TEAL, font="Consolas")
        add_text(slide, cl + Inches(0.2), ct + Inches(0.65),
                 card_w - Inches(0.4), card_h - Inches(0.85),
                 body, size=11.5, color=WHITE_85)

    # Right: AIOps stack (compact)
    aios_left = Inches(6.5)
    aios_w = Inches(6.2)
    aios_top = grid_top
    aios_h = Inches(4.15)

    # WVI envelope
    env = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  aios_left, aios_top, aios_w, aios_h)
    env.adjustments[0] = 0.05
    env.line.color.rgb = IQ_TEAL
    env.line.width = Pt(1.5)
    env.line.dash_style = 7
    env.fill.solid()
    env.fill.fore_color.rgb = IQ_TEAL
    env.fill.transparency = 0.94

    # AIOps tag
    add_pill(slide, aios_left + Inches(0.2), aios_top + Inches(0.1),
             Inches(1.6), Inches(0.3), "WVI · AIOps", IQ_TEAL, size=9)

    # Stack inside
    sx = aios_left + Inches(0.3)
    sw = aios_w - Inches(0.6)
    sy = aios_top + Inches(0.55)
    half = (sw - Inches(0.15)) / 2

    # Issue prevention band
    b = add_solid(slide, sx, sy, sw, Inches(0.42), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.6
    add_text(slide, sx, sy + Inches(0.06), sw, Inches(0.32),
             "Proactive + reactive issue prevention", size=11.5, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.52)

    # Two agents
    cd = add_solid(slide, sx, sy, half, Inches(0.55), AZURE_BLUE, line=AZURE_BLUE)
    cd.fill.transparency = 0.78
    add_text(slide, sx, sy + Inches(0.05), half, Inches(0.25),
             "Config Drift Agent", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sx, sy + Inches(0.3), half, Inches(0.22),
             "PROACTIVE", size=8.5, color=AZURE_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    rc = add_solid(slide, sx + half + Inches(0.15), sy, half, Inches(0.55),
                   ORACLE_RED, line=ORACLE_RED)
    rc.fill.transparency = 0.78
    add_text(slide, sx + half + Inches(0.15), sy + Inches(0.05), half, Inches(0.25),
             "Troubleshooting + RCA", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sx + half + Inches(0.15), sy + Inches(0.3), half, Inches(0.22),
             "REACTIVE", size=8.5, color=ORACLE_RED, align=PP_ALIGN.CENTER, font="Consolas")
    sy += Inches(0.65)

    # Common Agent Platform
    b = add_solid(slide, sx, sy, sw, Inches(0.42), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.55
    add_text(slide, sx, sy + Inches(0.06), sw, Inches(0.32),
             "Common Agent Platform", size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.52)

    # 3 pillars
    pillars = ["Configuration Validation", "Workload Health", "Contextual Support"]
    pw = (sw - Inches(0.3)) / 3
    for j, p in enumerate(pillars):
        px = sx + j * (pw + Inches(0.15))
        bx = add_solid(slide, px, sy, pw, Inches(0.48), WHITE, line=MUTED)
        bx.fill.transparency = 0.88
        add_text(slide, px, sy + Inches(0.06), pw, Inches(0.18),
                 p, size=9, color=AZURE_BLUE, align=PP_ALIGN.CENTER, bold=True)
    sy += Inches(0.56)

    # Foundation
    b = add_solid(slide, sx, sy, sw, Inches(0.42), IQ_TEAL, line=IQ_TEAL)
    b.fill.transparency = 0.55
    add_text(slide, sx, sy + Inches(0.06), sw, Inches(0.32),
             "Quality · Resiliency · Security",
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom takeaway band
    add_text(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.9),
             "Drift caught before incidents. RCA at machine speed. WVI-level rollups replace dashboard sprawl. The on-call rotation is quieter — and what fires is grounded.",
             size=12.5, color=WHITE_85)


# ---------------------------------------------------------------------------
# Slide 6 — The unified outcome (4 stakeholder lenses)
# ---------------------------------------------------------------------------
OUTCOMES = [
    {
        "label": "For the knowledge worker",
        "who": "Operators · planners · finance · supply-chain · HR · sales",
        "body": "Agents draft the email, the Teams message, the PO change. Decision in the calendar before the meeting. Hours back per worker per week.",
        "color": IQ_YELLOW,
    },
    {
        "label": "For the infra operator",
        "who": "SREs · platform engineers · ops leads · security",
        "body": "Drift caught before incidents. RCA at machine speed. WVI rollups replace dashboard sprawl. Quieter on-call. Grounded incidents.",
        "color": AZURE_BLUE,
    },
    {
        "label": "For the business leader",
        "who": "CEOs · COOs · CIOs · BU leaders",
        "body": "A frontier firm: agent-augmented, data-grounded, cross-stack, governed by default. Measurable productivity, cost-to-serve, risk reduction.",
        "color": IQ_TEAL,
    },
    {
        "label": "For the regulator / auditor",
        "who": "Risk · compliance · audit · sovereign authorities",
        "body": "End-to-end lineage from Oracle source to Copilot answer. Data Safe → Purview. Every agent action carries citation + policy evidence.",
        "color": ORACLE_RED,
    },
]


def outcomes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE UNIFIED OUTCOME",
        size=11, color=IQ_YELLOW, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "What the frontier firm looks like — across four stakeholder lenses.",
        size=24, bold=True, color=WHITE,
    )

    # 2x2 grid
    grid_top = Inches(1.85)
    card_w = Inches(6.0)
    card_h = Inches(2.55)
    gap = Inches(0.2)
    left = Inches(0.6)

    for i, o in enumerate(OUTCOMES):
        col = i % 2
        row = i // 2
        cl = left + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)
        color = o["color"]
        add_card(slide, cl, ct, card_w, card_h, border=color, fill=NAVY_700)
        add_text(slide, cl + Inches(0.25), ct + Inches(0.22),
                 card_w - Inches(0.5), Inches(0.35),
                 o["label"].upper(), size=10.5, bold=True, color=color, font="Consolas")
        add_text(slide, cl + Inches(0.25), ct + Inches(0.6),
                 card_w - Inches(0.5), Inches(0.35),
                 o["who"], size=11, color=MUTED, font="Consolas")
        # Divider
        line = slide.shapes.add_connector(1, cl + Inches(0.25), ct + Inches(1.0),
                                          cl + card_w - Inches(0.25), ct + Inches(1.0))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        # Body
        add_text(slide, cl + Inches(0.25), ct + Inches(1.15),
                 card_w - Inches(0.5), card_h - Inches(1.35),
                 o["body"], size=13, color=WHITE_85)


# ---------------------------------------------------------------------------
# Slide 7 — The ask
# ---------------------------------------------------------------------------
ASK_COMMITS_DECK = [
    ("Commit · Data estate landing zones",
     "Productised OD@A landing zones per industry. Joint Oracle Cloud Lift + Microsoft funding for the first three lighthouse customers per vertical."),
    ("Commit · Microsoft Agent Framework adoption",
     "Microsoft Agent Framework as the canonical build path; A2A + MCP as the canonical interop standards with Oracle Fusion AI Agents. Reference architectures co-published."),
    ("Commit · WVI + AIOps platform",
     "Workload Virtual Instance as a productised Azure construct. Common Agent Platform shipping Config Drift, Troubleshooting & RCA agents on a Quality / Resiliency / Security foundation."),
    ("Commit · Joint customer advisory board",
     "Standing CAB across all three workstreams. Quarterly cadence, 10+ named customers across FSI, Manufacturing, Retail, and Public Sector by end of Year 1."),
]


def ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE ASK  ·  FOUR COMMITS · ONE PER WORKSTREAM · ONE JOINT",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Four commits to ship this in the next 12 months.",
        size=24, bold=True, color=WHITE,
    )

    # 2x2 grid
    grid_top = Inches(1.85)
    card_w = Inches(6.0)
    card_h = Inches(2.2)
    gap = Inches(0.2)
    left = Inches(0.6)

    for i, (title, body) in enumerate(ASK_COMMITS_DECK):
        col = i % 2
        row = i // 2
        cl = left + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)
        add_card(slide, cl, ct, card_w, card_h, border=IQ_TEAL, fill=NAVY_700)
        add_text(slide, cl + Inches(0.25), ct + Inches(0.2),
                 card_w - Inches(0.5), Inches(0.35),
                 title.upper(), size=10.5, bold=True, color=IQ_TEAL, font="Consolas")
        add_text(slide, cl + Inches(0.25), ct + Inches(0.6),
                 card_w - Inches(0.5), card_h - Inches(0.8),
                 body, size=12.5, color=WHITE_85)

    # The bet closing band
    bet_top = Inches(6.6)
    bet = add_card(slide, Inches(0.6), bet_top, Inches(12.1), Inches(0.6),
                   border=IQ_TEAL, fill=NAVY_700, fill_alpha=0.4)
    bet.line.width = Pt(1.75)
    add_text(slide, Inches(0.85), bet_top + Inches(0.1),
             Inches(11.6), Inches(0.4),
             "THE BET  ·  Make Azure the destination Oracle workloads land on, the place their apps modernise on top of, and the surface their agents run from.",
             size=11.5, bold=True, color=IQ_TEAL, font="Consolas")


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    workstreams_slide(prs)
    complete_arch_slide(prs)
    interop_slide(prs)
    wvi_slide(prs)
    outcomes_slide(prs)
    ask_slide(prs)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "decks"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "frontier-transformation.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build()
