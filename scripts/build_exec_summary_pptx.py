"""Build the Exec Summary deck — the SLT-grade alliance briefing that
mirrors the /exec-summary portal page.

Output: docs/decks/exec-summary.pptx

Eight slides:
  1. Title — "Agentic AI is the next platform shift — Microsoft and Oracle
     can lead it together."
  2. Shift · Seam · Opportunity — 3 framing cards
  3. Market signals — 6 directional stat cards (33%, 15%, $450B+, 82%,
     88%, 40%+) + "why now" insight band
  4. The business case — 4 IQ tier columns (T0 -> T3) + compounding band
  5. Size of the prize — 4 sizing cards
  6. Joint consumption flywheel — 5-row table
  7. Investment & commercial frame — 3 phase cards
  8. The ask — 4 commits

Visual conventions, brand palette, and helper functions mirror
scripts/build_scenarios_pptx.py and scripts/build_oda_frontier_pptx.py
so all three decks read as a family.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Brand palette (matches the portal Tailwind config + sibling decks)
# ---------------------------------------------------------------------------
NAVY_900 = RGBColor(0x0B, 0x12, 0x20)
NAVY_700 = RGBColor(0x1A, 0x25, 0x40)
ORACLE_RED = RGBColor(0xC7, 0x46, 0x34)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
IQ_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
IQ_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_85 = RGBColor(0xD9, 0xDD, 0xE5)
WHITE_70 = RGBColor(0xB0, 0xB7, 0xC4)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

TIER_COLORS = {
    "T0": MUTED,
    "T1": AZURE_BLUE,
    "T2": IQ_YELLOW,
    "T3": IQ_TEAL,
}


# ---------------------------------------------------------------------------
# Helpers (mirror the sibling decks)
# ---------------------------------------------------------------------------
def add_background(slide, prs):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_900
    bg.shadow.inherit = False


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_pill(slide, left, top, width, height, text, color, *, size=9, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.line.color.rgb = color
    pill.line.width = Pt(0.75)
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.fill.transparency = 0.85
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Consolas"


def add_card(slide, left, top, width, height, *, border=MUTED, fill=NAVY_700, fill_alpha=0.6):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.line.color.rgb = border
    card.line.width = Pt(1)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.fill.transparency = fill_alpha
    card.shadow.inherit = False
    return card


def add_footer(slide, prs, text):
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.4),
        prs.slide_width, Inches(0.4),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.36), Inches(12), Inches(0.32),
        text, size=9, color=MUTED, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.6), Inches(12), Inches(0.4),
        "EXEC SUMMARY  ·  THE ENTERPRISE AGENT ALLIANCE  ·  ORACLE × MICROSOFT",
        size=11, color=IQ_YELLOW, font="Consolas",
    )

    add_text(
        slide, Inches(0.6), Inches(1.2), Inches(12), Inches(2.8),
        "Agentic AI is the next\nplatform shift —\nMicrosoft and Oracle\ncan lead it together.",
        size=40, bold=True, color=WHITE,
    )

    add_text(
        slide, Inches(0.6), Inches(5.1), Inches(12), Inches(1.5),
        "Oracle data. Microsoft agents. The complete agentic stack for the enterprise.",
        size=20, color=WHITE_85,
    )

    # Footer band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5),
        prs.slide_width, Inches(0.5),
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY_700
    add_text(
        slide, Inches(0.6), prs.slide_height - Inches(0.45), Inches(12), Inches(0.4),
        "Oracle anchors data of record where it's deployed.  ·  Microsoft extends the worker surface.  ·  Together: agentic transformation.",
        size=10, color=MUTED, font="Consolas",
    )


# ---------------------------------------------------------------------------
# Slide 2 — Shift · Seam · Opportunity (3 framing cards)
# ---------------------------------------------------------------------------
FRAMING_CARDS = [
    {
        "label": "THE SHIFT",
        "title": "From answering to acting",
        "body": "Agentic AI moves past chat. Agents reason over enterprise data, plan multi-step work, and execute with grounded citations — they draft the email, update the record, and close the loop. Every platform vendor is racing to be where these agents live.",
        "color": MUTED,
        "emphasised": False,
    },
    {
        "label": "THE SEAM",
        "title": "Where the work actually happens",
        "body": "Across the Fortune 2000, Oracle is one of the deepest stores of record-grade enterprise data — Fusion SaaS, Database@Azure, Oracle Cloud Infrastructure, on-prem databases. Microsoft 365 is where most information workers spend their day — Outlook, Teams, Copilot. Closing the seam is where agentic transformation actually happens.",
        "color": ORACLE_RED,
        "emphasised": False,
    },
    {
        "label": "THE OPPORTUNITY",
        "title": "Lead it together",
        "body": "Oracle's depth of business data plus Microsoft's agent runtime (Foundry), governance (Purview), and worker surface (Copilot, Teams) form the deepest enterprise agentic stack on offer. No other partnership covers data, reasoning, and the worker end-to-end at this scope.",
        "color": IQ_TEAL,
        "emphasised": True,
    },
]


def framing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE FRAME",
        size=11, color=MUTED, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.7),
        "The shift. The seam. The opportunity.",
        size=28, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4),
        "Why a joint Oracle + Microsoft motion is the only one that closes the agentic-AI seam for the Fortune 2000.",
        size=13, color=WHITE_70,
    )

    # 3 cards horizontally
    card_top = Inches(2.2)
    card_h = Inches(4.5)
    card_w = Inches(4.0)
    gap = Inches(0.15)
    left = Inches(0.6)

    for i, fc in enumerate(FRAMING_CARDS):
        cl = left + i * (card_w + gap)
        color = fc["color"]
        card = add_card(slide, cl, card_top, card_w, card_h,
                        border=color, fill=NAVY_700,
                        fill_alpha=0.35 if fc["emphasised"] else 0.65)
        if fc["emphasised"]:
            card.line.width = Pt(2.25)

        add_text(slide, cl + Inches(0.25), card_top + Inches(0.25),
                 card_w - Inches(0.5), Inches(0.3),
                 fc["label"], size=10, bold=True, color=color, font="Consolas")
        add_text(slide, cl + Inches(0.25), card_top + Inches(0.65),
                 card_w - Inches(0.5), Inches(0.85),
                 fc["title"], size=18, bold=True, color=WHITE)
        # Divider
        line = slide.shapes.add_connector(1, cl + Inches(0.25), card_top + Inches(1.65),
                                          cl + card_w - Inches(0.25), card_top + Inches(1.65))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        # Body
        add_text(slide, cl + Inches(0.25), card_top + Inches(1.85),
                 card_w - Inches(0.5), card_h - Inches(2.1),
                 fc["body"], size=12, color=WHITE_85)

    add_footer(slide=slide, prs=prs,
               text="Reference: /exec-summary  ·  The frame · §1")


# ---------------------------------------------------------------------------
# Slide 3 — Market signals (6 directional stats)
# ---------------------------------------------------------------------------
SIGNALS = [
    ("33%", "of enterprise software applications will include agentic AI by 2028 — up from <1% in 2024.", "Gartner, 2024", AZURE_BLUE, False),
    ("15%", "of day-to-day work decisions will be made autonomously by agentic AI by 2028 — up from 0% in 2024.", "Gartner, 2024", IQ_YELLOW, False),
    ("$450B+", "projected agentic enterprise software market by 2035 — ~30% of total enterprise app revenue.", "Gartner, 2025", IQ_TEAL, False),
    ("82%", "of leaders plan to use agents to expand workforce capacity in the next 12–18 months. 81% expect agents in operations.", "Microsoft Work Trend Index, 2025", MUTED, False),
    ("88%", "of organisations use AI in at least one function — but only 5.5% see real financial returns. The gap is data + integration.", "McKinsey State of AI, 2025", MUTED, False),
    ("40%+", "of agentic AI projects will be cancelled by end of 2027 — escalating costs, unclear value, weak governance. Architecture matters.", "Gartner, 2025", ORACLE_RED, True),
]


def signals_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "MARKET SIGNALS",
        size=11, color=AZURE_BLUE, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "The numbers say move now.",
        size=26, bold=True, color=WHITE,
    )

    # 2x3 grid of stat cards
    grid_top = Inches(1.7)
    card_w = Inches(4.0)
    card_h = Inches(2.0)
    gap_x = Inches(0.15)
    gap_y = Inches(0.18)
    left_start = Inches(0.6)

    for i, (value, body, source, color, warning) in enumerate(SIGNALS):
        col = i % 3
        row = i // 3
        cl = left_start + col * (card_w + gap_x)
        ct = grid_top + row * (card_h + gap_y)

        card = add_card(slide, cl, ct, card_w, card_h, border=color, fill=NAVY_700,
                        fill_alpha=0.5 if warning else 0.65)
        if warning:
            card.line.width = Pt(1.75)

        # Headline number
        add_text(slide, cl + Inches(0.25), ct + Inches(0.18),
                 card_w - Inches(0.5), Inches(0.65),
                 value, size=30, bold=True, color=WHITE)
        # Warning glyph for the 40% card
        if warning:
            add_text(slide, cl + card_w - Inches(0.75), ct + Inches(0.25),
                     Inches(0.5), Inches(0.4),
                     "⚠", size=20, color=ORACLE_RED, align=PP_ALIGN.RIGHT)
        # Body
        add_text(slide, cl + Inches(0.25), ct + Inches(0.85),
                 card_w - Inches(0.5), Inches(0.85),
                 body, size=11, color=WHITE_85)
        # Source
        add_text(slide, cl + Inches(0.25), ct + card_h - Inches(0.4),
                 card_w - Inches(0.5), Inches(0.28),
                 f"SOURCE · {source.upper()}",
                 size=9, color=color, font="Consolas", bold=True)

    # "Why now" insight band
    insight_top = grid_top + 2 * (card_h + gap_y) + Inches(0.1)
    insight_card = add_card(slide, Inches(0.6), insight_top, Inches(12.1), Inches(1.0),
                            border=IQ_YELLOW, fill=NAVY_700, fill_alpha=0.45)
    insight_card.line.width = Pt(1.5)
    add_text(slide, Inches(0.85), insight_top + Inches(0.12),
             Inches(11.6), Inches(0.3),
             "WHY NOW", size=10, bold=True, color=IQ_YELLOW, font="Consolas")
    add_text(slide, Inches(0.85), insight_top + Inches(0.4),
             Inches(11.6), Inches(0.6),
             "Adoption is mass-market, but value capture is rare. Bolt-on stacks stall; integrated ones compound. For Fortune 2000 enterprises running Oracle, the data-to-worker seam is among the largest unaddressed opportunities in enterprise software — and first-mover lead time compounds every quarter the alliance is shipping.",
             size=11.5, color=WHITE_85)


# ---------------------------------------------------------------------------
# Slide 4 — The business case (4 IQ tiers)
# ---------------------------------------------------------------------------
TIERS = [
    ("T0", "Baseline", "Lookup agent", "Oracle data",
     "Knows where the data lives. Describes tables, fields, definitions. Cannot interpret.",
     "“Here are the Oracle SCM tables that hold supplier and PO data.”",
     "Foundation only — table stakes. $0 incremental decision value."),
    ("T1", "+ Fabric IQ", "Insight agent", "Semantic",
     "Joins Oracle's semantic model. Quantifies, aggregates, trends — answers what the data says.",
     "“EMEA Q3 exposure is $47.2M across 84 suppliers; on-time delivery has fallen to 87%.”",
     "Quarterly review prep: days → minutes. Defensible numbers without SQL."),
    ("T2", "+ Foundry IQ", "Decision agent", "Reasoning",
     "Reasons across the data. Ranks moves, projects $-impact, cites every claim back to Oracle source.",
     "“5 prioritised de-risk plays. $18.6M mitigatable in Q3. First move: dual-source 25% of Adriatica volume.”",
     "Replaces an analyst week with a conversation. $-grade decisions on demand."),
    ("T3", "+ Work IQ", "Execution agent", "Execution",
     "Personalised to the worker via M365 Graph. Drafts the email, the Teams message, the meeting agenda — ready to send.",
     "“Drafted: green-light email to Maria, Teams ping to Tomás, pre-filled 10am ops-sync agenda — all in your queue.”",
     "The decision is in the calendar before the meeting starts. Hours back per worker, per week."),
]


def business_case_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE BUSINESS CASE · IQ TIER PROGRESSION",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Each IQ layer turns Oracle data into a more valuable kind of agent.",
        size=23, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.4),
        "Same Oracle Fusion source. Four agent archetypes. The value isn't additive — it's multiplicative.",
        size=12, color=WHITE_70,
    )

    # 4 tier columns
    col_top = Inches(2.05)
    col_h = Inches(4.5)
    col_w = Inches(2.95)
    gap = Inches(0.15)
    left_start = Inches(0.6)

    for i, (tier_id, label, archetype, tag, blurb, example, value) in enumerate(TIERS):
        cl = left_start + i * (col_w + gap)
        color = TIER_COLORS[tier_id]

        add_card(slide, cl, col_top, col_w, col_h, border=color, fill=NAVY_700)

        # Tier badge + tag chip on top
        add_text(slide, cl + Inches(0.2), col_top + Inches(0.18),
                 Inches(1.4), Inches(0.3),
                 f"{tier_id}  ·  {label}", size=10, bold=True, color=color, font="Consolas")
        add_pill(slide, cl + col_w - Inches(1.05), col_top + Inches(0.18),
                 Inches(0.85), Inches(0.28), tag, color, size=8)

        # Archetype
        add_text(slide, cl + Inches(0.2), col_top + Inches(0.6),
                 col_w - Inches(0.4), Inches(0.4),
                 archetype, size=15, bold=True, color=WHITE)

        # Blurb
        add_text(slide, cl + Inches(0.2), col_top + Inches(1.0),
                 col_w - Inches(0.4), Inches(1.0),
                 blurb, size=11, color=WHITE_85)

        # Example label
        add_text(slide, cl + Inches(0.2), col_top + Inches(2.05),
                 col_w - Inches(0.4), Inches(0.25),
                 "EXAMPLE", size=9, color=MUTED, font="Consolas")
        add_text(slide, cl + Inches(0.2), col_top + Inches(2.3),
                 col_w - Inches(0.4), Inches(1.3),
                 example, size=10.5, color=WHITE_85)

        # Value label + body
        add_text(slide, cl + Inches(0.2), col_top + Inches(3.5),
                 col_w - Inches(0.4), Inches(0.25),
                 "VALUE", size=9, color=MUTED, font="Consolas")
        add_text(slide, cl + Inches(0.2), col_top + Inches(3.75),
                 col_w - Inches(0.4), Inches(0.7),
                 value, size=10.5, color=color, bold=True)

    # Compounding band
    band_top = Inches(6.75)
    band = add_card(slide, Inches(0.6), band_top, Inches(12.1), Inches(0.55),
                    border=IQ_TEAL, fill=NAVY_700, fill_alpha=0.45)
    band.line.width = Pt(1.5)
    add_text(slide, Inches(0.85), band_top + Inches(0.1),
             Inches(11.6), Inches(0.4),
             "THE COMPOUNDING EFFECT  ·  Most agentic projects ship at T1 and stall — they report but never act. The alliance ships the full T0→T3 progression. By Tier 3 the agent isn't reporting on the business — it's running parts of it.",
             size=11, color=IQ_TEAL, font="Consolas", bold=True)


# ---------------------------------------------------------------------------
# Slide 5 — Size of the prize (4 stat cards)
# ---------------------------------------------------------------------------
PRIZE_STATS = [
    ("$450B+", "Total agentic TAM by 2035", "Alliance plays for outsized share · Gartner", IQ_YELLOW, True),
    ("~6,000", "Customers in scope", "Estimated Fusion + M365 overlap; under joint validation", MUTED, False),
    ("9", "Reusable building blocks", "3 patterns + 2 accelerators + 4 industry blueprints", MUTED, False),
    ("~12 wks", "Target time to first T2 agent", "With accelerators GA; bespoke today: 9 – 12 months", IQ_TEAL, True),
]


def prize_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "SIZE OF THE PRIZE",
        size=11, color=IQ_YELLOW, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "The opportunity sized for the alliance.",
        size=26, bold=True, color=WHITE,
    )

    # 4 stat cards in a row
    card_top = Inches(2.0)
    card_h = Inches(3.0)
    card_w = Inches(2.95)
    gap = Inches(0.15)
    left_start = Inches(0.6)

    for i, (value, label, sub, color, emphasised) in enumerate(PRIZE_STATS):
        cl = left_start + i * (card_w + gap)
        card = add_card(slide, cl, card_top, card_w, card_h, border=color, fill=NAVY_700,
                        fill_alpha=0.4 if emphasised else 0.6)
        if emphasised:
            card.line.width = Pt(1.75)
        # Label
        add_text(slide, cl + Inches(0.2), card_top + Inches(0.25),
                 card_w - Inches(0.4), Inches(0.3),
                 label.upper(), size=10, bold=True, color=color, font="Consolas")
        # Value
        add_text(slide, cl + Inches(0.2), card_top + Inches(0.7),
                 card_w - Inches(0.4), Inches(1.2),
                 value, size=42, bold=True, color=WHITE)
        # Sub
        add_text(slide, cl + Inches(0.2), card_top + Inches(2.0),
                 card_w - Inches(0.4), Inches(0.9),
                 sub, size=11, color=WHITE_85)

    # Closing paragraph
    add_text(
        slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5),
        "The prize isn't just license revenue — it's consumption across both stacks. Every agent invocation drives Microsoft Fabric capacity, Foundry tokens, and Copilot seats — while compounding Oracle Database@Azure, OCI, and Fusion subscription consumption. The flywheel on the next slide shows the mechanics.",
        size=13, color=WHITE_85,
    )

    add_footer(slide=slide, prs=prs,
               text="Reference: /exec-summary  ·  Size of the prize · §4")


# ---------------------------------------------------------------------------
# Slide 6 — Joint consumption flywheel (5-row table)
# ---------------------------------------------------------------------------
FLYWHEEL_ROWS = [
    ("1", "Worker asks question in Copilot or Fusion",
     ["Copilot for M365 (per-seat ARR)"], ["Oracle Fusion (per-user subscription)"]),
    ("2", "Foundry agent reasons, retrieves, plans",
     ["Azure AI Foundry (token consumption)", "Microsoft Agent 365 (control plane)"],
     ["Fusion AI Agents (when invoked)"]),
    ("3", "Reads from Fabric semantic model",
     ["Microsoft Fabric capacity (CU)", "OneLake storage"],
     ["Oracle GoldenGate (replication metering)"]),
    ("4", "Source data lives in Oracle (Fusion / DB@Azure / OCI / on-prem)",
     ["Database@Azure passthrough"],
     ["Database@Azure consumption", "OCI compute & storage", "Fusion Apps (subscription)"]),
    ("5", "Governance, lineage, and audit flow back",
     ["Microsoft Purview"],
     ["Oracle Data Safe", "Oracle Audit Vault"]),
]


def flywheel_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "JOINT CONSUMPTION FLYWHEEL",
        size=11, color=IQ_YELLOW, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "Every agent invocation drives revenue across both clouds.",
        size=24, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.4),
        "One workflow. Two stacks. Five revenue lines simultaneously.",
        size=12, color=WHITE_70,
    )

    # Header
    header_top = Inches(2.0)
    col_x = [Inches(0.6), Inches(1.45), Inches(5.6), Inches(9.4)]
    col_w = [Inches(0.75), Inches(4.0), Inches(3.7), Inches(3.3)]

    headers = [
        ("STEP", MUTED),
        ("WORKER ACTION", MUTED),
        ("MICROSOFT CONSUMPTION", AZURE_BLUE),
        ("ORACLE CONSUMPTION", ORACLE_RED),
    ]
    for (label, color), x, w in zip(headers, col_x, col_w):
        add_text(slide, x, header_top, w, Inches(0.3), label,
                 size=10, color=color, font="Consolas", bold=True)

    underline = slide.shapes.add_connector(1, Inches(0.6), header_top + Inches(0.36),
                                           Inches(12.7), header_top + Inches(0.36))
    underline.line.color.rgb = MUTED
    underline.line.width = Pt(0.5)
    underline.line.fill.transparency = 0.4

    # Rows
    row_top = Inches(2.55)
    row_h = Inches(0.78)
    for i, (step, action, ms_list, ora_list) in enumerate(FLYWHEEL_ROWS):
        y = row_top + i * row_h
        # Alternating band
        if i % 2 == 0:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y,
                                          Inches(12.1), row_h - Inches(0.05))
            band.line.fill.background()
            band.fill.solid()
            band.fill.fore_color.rgb = NAVY_700
            band.fill.transparency = 0.55

        # Step circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        col_x[0] + Inches(0.05), y + Inches(0.18),
                                        Inches(0.42), Inches(0.42))
        circle.line.color.rgb = IQ_YELLOW
        circle.line.width = Pt(1)
        circle.fill.solid()
        circle.fill.fore_color.rgb = IQ_YELLOW
        circle.fill.transparency = 0.85
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = ctf.margin_right = 0
        ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        crun = cp.add_run()
        crun.text = step
        crun.font.size = Pt(12)
        crun.font.bold = True
        crun.font.color.rgb = IQ_YELLOW
        crun.font.name = "Consolas"

        # Worker action
        add_text(slide, col_x[1], y + Inches(0.18), col_w[1], Inches(0.55),
                 action, size=12, color=WHITE, bold=True)

        # Microsoft chips
        ms_x = col_x[2]
        ms_y = y + Inches(0.2)
        for label in ms_list:
            # Estimate pill width from text length
            w = Inches(max(1.0, min(3.4, 0.08 * len(label) + 0.4)))
            add_pill(slide, ms_x, ms_y, w, Inches(0.32), label, AZURE_BLUE, size=8.5)
            ms_y += Inches(0.38)

        # Oracle chips
        ora_x = col_x[3]
        ora_y = y + Inches(0.2)
        for label in ora_list:
            w = Inches(max(1.0, min(3.0, 0.08 * len(label) + 0.4)))
            add_pill(slide, ora_x, ora_y, w, Inches(0.32), label, ORACLE_RED, size=8.5)
            ora_y += Inches(0.38)

    # Footer paragraph
    foot_y = row_top + 5 * row_h + Inches(0.1)
    add_text(
        slide, Inches(0.6), foot_y, Inches(12.1), Inches(0.85),
        "The flywheel: more Oracle data into OneLake → more Fabric capacity consumed; more agents in Foundry → more tokens; more Copilot invocations → more pull on both. Microsoft monetises Fabric + Foundry + Copilot simultaneously; Oracle monetises Database@Azure + OCI + Fusion in lockstep.",
        size=11.5, color=WHITE_85,
    )

    add_footer(slide=slide, prs=prs,
               text="Reference: /exec-summary  ·  Joint consumption flywheel · §5")


# ---------------------------------------------------------------------------
# Slide 7 — Investment & commercial frame
# ---------------------------------------------------------------------------
COMMERCIAL_PHASES = [
    {
        "label": "YEAR 1",
        "title": "Co-sell, separate billing",
        "body": "Lowest-friction path. Both sales orgs co-sell with shared pipeline credit; each vendor bills for its own consumption. Joint reference architecture and shared accelerators in market.",
        "color": MUTED,
    },
    {
        "label": "YEAR 2",
        "title": "Joint SKU + marketplace",
        "body": "Single alliance SKU on Azure Marketplace and Oracle Cloud Marketplace; unified GTM with named field motion; co-marketed in industry events; ISV / SI program announced jointly.",
        "color": AZURE_BLUE,
    },
    {
        "label": "CO-FUNDED BUILD",
        "title": "Accelerators & blueprints",
        "body": "Ontology Bridge, Governance Bridge, and the four industry blueprints co-developed and co-owned. Shared engineering investment; shared roadmap; shared customer advisory board.",
        "color": IQ_YELLOW,
    },
]


def commercial_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "INVESTMENT & COMMERCIAL FRAME",
        size=11, color=AZURE_BLUE, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "How the alliance gets paid for.",
        size=26, bold=True, color=WHITE,
    )
    add_text(
        slide, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.4),
        "Crawl-walk-run on the commercial structure. Engineering investment is co-funded from day one.",
        size=12, color=WHITE_70,
    )

    # 3 phase cards
    card_top = Inches(2.1)
    card_h = Inches(4.6)
    card_w = Inches(4.0)
    gap = Inches(0.15)
    left = Inches(0.6)

    for i, phase in enumerate(COMMERCIAL_PHASES):
        cl = left + i * (card_w + gap)
        color = phase["color"]
        add_card(slide, cl, card_top, card_w, card_h, border=color, fill=NAVY_700)
        # Big label
        add_text(slide, cl + Inches(0.25), card_top + Inches(0.3),
                 card_w - Inches(0.5), Inches(0.4),
                 phase["label"], size=11, bold=True, color=color, font="Consolas")
        # Title
        add_text(slide, cl + Inches(0.25), card_top + Inches(0.8),
                 card_w - Inches(0.5), Inches(1.0),
                 phase["title"], size=20, bold=True, color=WHITE)
        # Divider
        line = slide.shapes.add_connector(1, cl + Inches(0.25), card_top + Inches(2.05),
                                          cl + card_w - Inches(0.25), card_top + Inches(2.05))
        line.line.color.rgb = color
        line.line.width = Pt(0.5)
        line.line.fill.transparency = 0.6
        # Body
        add_text(slide, cl + Inches(0.25), card_top + Inches(2.25),
                 card_w - Inches(0.5), card_h - Inches(2.5),
                 phase["body"], size=12, color=WHITE_85)


# ---------------------------------------------------------------------------
# Slide 8 — The ask (4 commits)
# ---------------------------------------------------------------------------
COMMITS = [
    ("Commit 1 · Joint sponsorship",
     "Both leadership teams publicly endorse the Enterprise Agent Alliance. Joint announcement at the next major event from each side."),
    ("Commit 2 · Co-funded GTM",
     "Named field organisation with shared targets and pipeline credit. Co-marketing budget agreed for Year 1 with quarterly review cadence."),
    ("Commit 3 · Accelerators GA",
     "Ontology Bridge and Governance Bridge accelerators committed to GA by end of Phase 2. Co-owned engineering, co-owned roadmap."),
    ("Commit 4 · Joint Customer Advisory Board",
     "Standing CAB with quarterly cadence and 10+ named participants by end of Year 1. All four industry blueprints co-developed by end of Phase 3; global SIs brought in under a co-funded enablement programme."),
]


def ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, prs)

    add_text(
        slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.35),
        "THE ASK  ·  DECISION-GRADE",
        size=11, color=IQ_TEAL, font="Consolas",
    )
    add_text(
        slide, Inches(0.6), Inches(0.8), Inches(12), Inches(0.6),
        "What we're asking the SLT to commit to.",
        size=26, bold=True, color=WHITE,
    )

    # 2x2 grid of commits
    grid_top = Inches(1.85)
    card_w = Inches(6.0)
    card_h = Inches(2.3)
    gap = Inches(0.2)
    left_start = Inches(0.6)

    for i, (title, body) in enumerate(COMMITS):
        col = i % 2
        row = i // 2
        cl = left_start + col * (card_w + gap)
        ct = grid_top + row * (card_h + gap)

        add_card(slide, cl, ct, card_w, card_h, border=IQ_TEAL, fill=NAVY_700)
        # Commit label
        add_text(slide, cl + Inches(0.25), ct + Inches(0.2),
                 card_w - Inches(0.5), Inches(0.35),
                 title.upper(), size=10.5, bold=True, color=IQ_TEAL, font="Consolas")
        # Body
        add_text(slide, cl + Inches(0.25), ct + Inches(0.65),
                 card_w - Inches(0.5), card_h - Inches(0.85),
                 body, size=13, color=WHITE_85)

    # Closing band
    bet_top = Inches(6.65)
    bet_card = add_card(slide, Inches(0.6), bet_top, Inches(12.1), Inches(0.55),
                        border=IQ_TEAL, fill=NAVY_700, fill_alpha=0.4)
    bet_card.line.width = Pt(1.75)
    add_text(slide, Inches(0.85), bet_top + Inches(0.1),
             Inches(11.6), Inches(0.4),
             "FIRST-MOVER LEAD TIME COMPOUNDS  ·  Every quarter the alliance is shipping is a quarter competitors fall further behind on the data + worker seam.",
             size=11.5, bold=True, color=IQ_TEAL, font="Consolas")


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    framing_slide(prs)
    signals_slide(prs)
    business_case_slide(prs)
    prize_slide(prs)
    flywheel_slide(prs)
    commercial_slide(prs)
    ask_slide(prs)

    out_dir = Path(__file__).resolve().parent.parent / "docs" / "decks"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "exec-summary.pptx"
    prs.save(out_path)
    print(f"Wrote {out_path}")


if __name__ == "__main__":
    build()
